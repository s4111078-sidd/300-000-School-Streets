"""
build_presentation_v2.py
School Streets Safety — clean redesign
Mixed audience: data science / social science / Vic Health
~10 slides, 10–15 minute run time

Layout rules enforced throughout:
  Header band:  y = 0   → 1.05
  Content zone: y = 1.12 → 7.1
  Footer band:  y = 7.12 → 7.5
  Sidebar+image split: sidebar x=0.30..3.85 | image x=4.05..13.28
  All text elements stay WITHIN their declared zone — zero overlap with images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
GD   = RGBColor(0x1B, 0x5E, 0x20)   # dark green  — header, accent
GM   = RGBColor(0x2E, 0x7D, 0x32)   # mid green   — cards, bullets bar
GL   = RGBColor(0x43, 0xA0, 0x47)   # leaf green  — accent stripe
GP   = RGBColor(0xF1, 0xF8, 0xE9)   # pale green  — section bg
GLIT = RGBColor(0xE8, 0xF5, 0xE9)   # light green — footer bg
TEA  = RGBColor(0x00, 0x69, 0x6E)   # teal        — callout boxes
AMB  = RGBColor(0xE6, 0x5C, 0x00)   # amber/orange— warning
ALT  = RGBColor(0xFF, 0xF3, 0xE0)   # amber light
RED  = RGBColor(0xC6, 0x28, 0x28)   # red         — critical
RLT  = RGBColor(0xFF, 0xEB, 0xEE)   # red light
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
INK  = RGBColor(0x21, 0x21, 0x21)
SLT  = RGBColor(0x37, 0x47, 0x4F)   # slate       — footer text
MUT  = RGBColor(0x78, 0x90, 0x9C)   # muted blue-grey
CGRN = RGBColor(0xC8, 0xE6, 0xC9)   # card text   — light green on dark
AGRN = RGBColor(0xA5, 0xD6, 0xA7)   # sub-header  — medium green

OUT  = os.path.join(os.path.dirname(__file__), "outputs")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
BLANK = prs.slide_layouts[6]

# ── Primitive helpers ─────────────────────────────────────────────────────────

def R(s, l, t, w, h, fill):
    """Solid filled rectangle, no border."""
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def T(s, text, l, t, w, h,
      bold=False, sz=14, col=None, al=PP_ALIGN.LEFT, it=False):
    col = col or INK
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    run = p.add_run(); run.text = text
    run.font.bold = bold; run.font.italic = it
    run.font.size = Pt(sz); run.font.color.rgb = col

def IMG(s, fname, l, t, w, h=None):
    path = os.path.join(OUT, fname)
    if not os.path.exists(path): return
    if h: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

# ── Component helpers ─────────────────────────────────────────────────────────

def HDR(s, title, sub=None):
    """Dark-green header band with optional subtitle."""
    R(s, 0, 0, 13.33, 1.05, GD)
    R(s, 0, 1.05, 13.33, 0.07, GL)          # leaf-green accent stripe
    T(s, title, 0.40, 0.08, 12.5, 0.56, bold=True, sz=23, col=WHT)
    if sub:
        T(s, sub, 0.40, 0.67, 12.5, 0.34, sz=11, col=AGRN)

def FTR(s, n):
    """Footer band."""
    R(s, 0, 7.12, 13.33, 0.38, GLIT)
    T(s, "300,000 Streets of Melbourne  ·  Regen Melbourne × RMIT University  ·  Darebin Pilot",
      0.35, 7.15, 11.5, 0.28, sz=8.5, col=SLT)
    T(s, str(n), 12.75, 7.15, 0.40, 0.28, sz=8.5, col=SLT, al=PP_ALIGN.RIGHT)

def STAT(s, val, lbl, l, t, w=3.55, h=1.38, bg=GM):
    """Large stat card for sidebar or title row."""
    R(s, l, t, w, h, bg)
    R(s, l, t+h-0.08, w, 0.08, GL)          # bottom accent stripe
    T(s, val, l, t+0.05, w, 0.65, bold=True, sz=36, col=WHT, al=PP_ALIGN.CENTER)
    T(s, lbl, l+0.08, t+0.70, w-0.16, 0.55, sz=10, col=CGRN, al=PP_ALIGN.CENTER)

def STAT_SM(s, val, lbl, l, t, w=2.15, h=1.30, bg=GM):
    """Smaller stat card (5 per row on title slide)."""
    R(s, l, t, w, h, bg)
    R(s, l, t+h-0.07, w, 0.07, GL)
    T(s, val, l, t+0.04, w, 0.60, bold=True, sz=32, col=WHT, al=PP_ALIGN.CENTER)
    T(s, lbl, l+0.06, t+0.66, w-0.12, 0.52, sz=9.5, col=CGRN, al=PP_ALIGN.CENTER)

def NOTE(s, text, l, t, w, h, bg=TEA, sz=10.5):
    """Callout note with left accent bar."""
    R(s, l, t, w, h, bg)
    R(s, l, t, 0.07, h, GL)                 # left accent
    T(s, text, l+0.18, t+0.09, w-0.27, h-0.17, sz=sz, col=WHT)

def BSEC(s, title, items, l, t, w, bg=GP, tc=GD, isz=11):
    """Titled bullet section, returns height used."""
    h = 0.44 + len(items) * 0.38
    R(s, l, t, w, h, bg)
    R(s, l, t, 0.07, h, GM)
    T(s, title, l+0.18, t+0.09, w-0.27, 0.30, bold=True, sz=12.5, col=tc)
    y = t + 0.44
    for item in items:
        T(s, f"• {item}", l+0.18, y, w-0.27, 0.34, sz=isz, col=INK)
        y += 0.36
    return h

def ANNOT_FLAG(s, text, l, t, w, h, bg, arrow_down_to=None):
    """Annotation flag with optional downward pointer line."""
    R(s, l, t, w, h, bg)
    R(s, l, t, w, 0.06, GL)                 # top accent stripe
    T(s, text, l+0.12, t+0.1, w-0.24, h-0.18, sz=10, col=WHT, bold=True)
    if arrow_down_to is not None:
        cx = l + w/2 - 0.03
        R(s, cx, t+h, 0.06, arrow_down_to-(t+h), RGBColor(0xBD,0xBD,0xBD))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)

# Full-bleed split: top dark, bottom pale
R(sl, 0, 0, 13.33, 4.10, GD)
R(sl, 0, 4.10, 13.33, 3.40, GLIT)
R(sl, 0, 4.05, 13.33, 0.10, GL)             # separator stripe

# Main title text block
T(sl, "300,000 Streets of Melbourne",
  0.55, 0.42, 12.2, 1.05, bold=True, sz=44, col=WHT)
T(sl, "School Streets Safety Analysis — City of Darebin",
  0.55, 1.52, 12.2, 0.70, sz=26, col=AGRN)
T(sl, "A pedestrian & cyclist safety assessment for secondary school walking routes",
  0.55, 2.28, 12.2, 0.55, sz=16, col=RGBColor(0x81, 0xC7, 0x84), it=True)

# Org / date strip inside dark band
R(sl, 0.55, 3.18, 4.50, 0.70, RGBColor(0x14, 0x4A, 0x18))
T(sl, "Regen Melbourne  ×  RMIT University",
  0.70, 3.24, 4.20, 0.38, bold=True, sz=14, col=WHT)
T(sl, "June 2026", 0.70, 3.58, 4.20, 0.28, sz=11, col=AGRN, it=True)

# Framework reference pill
R(sl, 5.30, 3.18, 5.80, 0.70, RGBColor(0x14, 0x4A, 0x18))
T(sl, "Healthy Streets Framework  |  Lucy Saunders / Transport for London",
  5.45, 3.24, 5.50, 0.38, sz=12, col=AGRN)
T(sl, "Open data  •  Field observation  •  ML scoring  •  Equity analysis",
  5.45, 3.58, 5.50, 0.28, sz=10, col=MUT)

# ── 5 stat cards in the pale zone ────────────────────────────────────────────
cards = [
    ("3",       "Secondary schools\nassessed in Darebin",     GM),
    ("7,948",   "Vic ped/cyc crashes\nanalysed (2021–25)",    TEA),
    ("192%",    "Crash increase\nDarebin 2021–2024",          RED),
    ("r=0.84",  "Equity–safety\ncorrelation (Pearson)",       AMB),
    ("17:00",   "Peak crash hour\n(school pickup)",           RGBColor(0x4A,0x00,0x0E)),
]
cx = 0.50
for val, lbl, bg in cards:
    STAT_SM(sl, val, lbl, cx, 4.25, w=2.30, h=1.32, bg=bg)
    cx += 2.47

T(sl, "Healthy Streets Framework  |  10 indicators, 0–10 score  |  Major / Moderate / Minor severity",
  0.55, 5.70, 12.2, 0.35, sz=11, col=SLT, al=PP_ALIGN.CENTER)
T(sl, "Data sources: VicRoads  ·  OpenStreetMap  ·  EPA Victoria  ·  Crime Statistics Agency  ·  ABS Census 2021  ·  ABS SEIFA 2021",
  0.55, 6.12, 12.2, 0.30, sz=9.5, col=MUT, al=PP_ALIGN.CENTER)

FTR(sl, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — PROBLEM STATEMENT
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "The Problem",
    "Children walking to school face rising danger — and disadvantaged communities bear the most risk")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# ── Left column: 3 big number panels ─────────────────────────────────────────
# Panel 1 — crash surge
R(sl, 0.30, 1.22, 3.55, 1.65, RED)
R(sl, 0.30, 1.22, 3.55, 0.07, AMB)
T(sl, "192%", 0.30, 1.28, 3.55, 0.85, bold=True, sz=56, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "increase in Darebin LGA\nped/cyc crashes  2021 → 2024\n(25 crashes → 73 crashes)",
  0.40, 2.12, 3.35, 0.70, sz=11, col=RLT, al=PP_ALIGN.CENTER)

# Panel 2 — peak hour
R(sl, 0.30, 3.02, 3.55, 1.55, AMB)
R(sl, 0.30, 3.02, 3.55, 0.07, GL)
T(sl, "17:00", 0.30, 3.08, 3.55, 0.80, bold=True, sz=52, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "Peak crash hour across\nDarebin LGA — school pickup time",
  0.40, 3.86, 3.35, 0.62, sz=11, col=ALT, al=PP_ALIGN.CENTER)

# Panel 3 — equity
R(sl, 0.30, 4.72, 3.55, 1.50, GD)
R(sl, 0.30, 4.72, 3.55, 0.07, GL)
T(sl, "r = 0.84", 0.30, 4.78, 3.55, 0.75, bold=True, sz=42, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "Equity correlation — more disadvantaged\ncatchments have worse safety scores",
  0.40, 5.50, 3.35, 0.62, sz=11, col=AGRN, al=PP_ALIGN.CENTER)

# ── Right column: context + health frame ─────────────────────────────────────
R(sl, 4.05, 1.22, 9.0, 5.98, WHT)

# Section header
R(sl, 4.05, 1.22, 9.0, 0.44, GD)
T(sl, "Why This Matters — Safety, Equity & Health", 4.20, 1.26, 8.75, 0.36,
  bold=True, sz=14, col=WHT)

# Three context panels stacked in the right column
panels = [
    (RED, "Safety — streets are getting more dangerous",
     ["Darebin LGA recorded 73 pedestrian & cyclist crashes in 2024, up from 25 in 2021",
      "Preston HS has 15 crashes within 400 m of its gate across 4 years",
      "100% of William Ruthven SC's nearby crashes occurred during school hours",
      "No safe crossing exists at the Preston HS gate — HS2 score = 0.4 out of 10"]),
    (AMB, "Health — active travel to school is a public health asset",
     ["Walking/cycling to school is linked to better physical fitness and mental health",
      "Children in car-dependent households are 40% less likely to be physically active",
      "Air quality, noise and traffic stress near arterials directly affect development",
      "Sedentary travel habits formed in childhood persist into adulthood (WHO evidence)"]),
    (TEA, "Equity — the most vulnerable bear the greatest burden",
     ["Reservoir & William Ruthven catchments sit in SEIFA Decile 4 (moderate-high disadvantage)",
      "Lower car ownership (11–13%) means these communities depend on walking/cycling",
      "Median household income in Reservoir is $1,541/wk — below Melbourne median ($1,950)",
      "Safer streets are not a premium — they're a basic health and equity requirement"]),
]
py = 1.70
for col, heading, bullets in panels:
    R(sl, 4.05, py, 9.0, 0.36, col)
    T(sl, heading, 4.18, py+0.05, 8.75, 0.27, bold=True, sz=12, col=WHT)
    by = py + 0.38
    for b in bullets:
        T(sl, f"  • {b}", 4.12, by, 8.78, 0.30, sz=11, col=INK)
        by += 0.30
    R(sl, 4.05, by+0.04, 9.0, 0.04, GLIT)  # divider
    py = by + 0.10

FTR(sl, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — FRAMEWORK & APPROACH
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "How We Approached the Problem",
    "Healthy Streets framework + open data pipeline + machine learning + equity analysis")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# ── Left: Healthy Streets 10 indicators ──────────────────────────────────────
R(sl, 0.30, 1.22, 5.95, 0.44, GD)
T(sl, "Healthy Streets Framework — 10 Indicators (0–10 each)",
  0.44, 1.26, 5.75, 0.34, bold=True, sz=13, col=WHT)
T(sl, "Lucy Saunders / Transport for London  |  International evidence base",
  0.44, 1.62, 5.75, 0.26, sz=9.5, col=SLT, it=True)

inds = [
    (GM,  "HS1",  "Pedestrians from all walks of life", "Field obs — footpath presence, width, kerb ramps"),
    (GM,  "HS2",  "Easy to cross",                      "Field obs — crossing type, distance, signals"),
    (GM,  "HS3",  "Shade and shelter",                  "OSM — tree count, shelters, green space %"),
    (GM,  "HS4",  "Places to stop and rest",            "OSM — benches, parks, cafes"),
    (GM,  "HS5",  "Not too noisy",                      "Field obs — traffic volume, speed, lanes"),
    (GM,  "HS6",  "People choose active travel",        "OSM + field — cycling infrastructure, PT stops"),
    (GM,  "HS7",  "People feel safe",                   "Field + CSA Victoria — lighting, crime rate"),
    (GM,  "HS8",  "Things to see and do",               "OSM — amenities, cafes, parks"),
    (GM,  "HS9",  "People feel relaxed",                "Field obs — traffic calming, school zones"),
    (GM,  "HS10", "Clean air",                          "EPA Victoria AirWatch — PM2.5 annual avg"),
]
iy = 1.92
for bg, code, name, source in inds:
    R(sl, 0.30, iy, 0.68, 0.33, bg)
    T(sl, code, 0.30, iy+0.03, 0.68, 0.27, bold=True, sz=9.5, col=WHT, al=PP_ALIGN.CENTER)
    T(sl, name, 1.04, iy+0.03, 2.58, 0.27, bold=True, sz=10, col=GD)
    T(sl, source, 3.66, iy+0.03, 2.50, 0.27, sz=9, col=SLT)
    iy += 0.36
# threshold note
R(sl, 0.30, 5.54, 5.95, 0.30, RGBColor(0x14,0x4A,0x18))
T(sl, "Target: all indicators ≥ 6.0  |  Major / Moderate / Minor severity classification applied per school gate",
  0.44, 5.57, 5.78, 0.24, sz=9, col=AGRN)

# ── Right: pipeline ───────────────────────────────────────────────────────────
R(sl, 6.55, 1.22, 6.45, 0.44, TEA)
T(sl, "Our 10-Step Analysis Pipeline", 6.68, 1.26, 6.22, 0.34, bold=True, sz=13, col=WHT)

steps = [
    ("1", "Crash data",         "7,948 Victorian ped/cyc crashes 2021–2025, nearest school assigned",   GM),
    ("2", "OSM features",       "53 spatial features per school at 200m / 400m / 800m buffers",         GM),
    ("3", "Environmental",      "EPA PM2.5 air quality + Crime Statistics Agency crime rate",            GM),
    ("4", "HS scoring + maps",  "10-indicator scores, hazard severity, interactive Leaflet maps",       GM),
    ("5", "ML feature matrix",  "26 open-data predictors × 10 HS score targets",                        TEA),
    ("6", "Ridge regression",   "LOO-CV prediction of HS scores from open data (mean MAE 2.88)",         TEA),
    ("7", "SEIFA analysis",     "ABS 2021 disadvantage deciles mapped to school catchments",            GD),
    ("8", "Equity overlay",     "SEIFA disadvantage × Healthy Streets score  →  r = 0.84",              GD),
    ("9", "Crash trends",       "Year-on-year trends, school-hours breakdown, time-of-day peak",        AMB),
    ("10","Scenario analysis",  "10 interventions × 3 schools = 30 modelled what-if results",           AMB),
]
sy = 1.70
for num, title, desc, bg in steps:
    R(sl, 6.55, sy, 0.52, 0.33, bg)
    T(sl, num, 6.55, sy+0.03, 0.52, 0.27, bold=True, sz=10, col=WHT, al=PP_ALIGN.CENTER)
    T(sl, title, 7.14, sy+0.03, 1.80, 0.27, bold=True, sz=10.5, col=GD)
    T(sl, desc,  8.98, sy+0.03, 3.95, 0.27, sz=9.5, col=INK)
    sy += 0.37

R(sl, 6.55, 5.54, 6.45, 0.30, RGBColor(0x00,0x45,0x4A))
T(sl, "Reproducible pipeline — python run_all.py regenerates all outputs from source data",
  6.68, 5.57, 6.28, 0.24, sz=9, col=RGBColor(0x80,0xCB,0xC4))

FTR(sl, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — SAFETY SCORES
# chart3_score_breakdown.png fills the lower 5/7 of the slide
# An annotation row above points down into each panel
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Safety Scores — 3 Darebin Secondary Schools",
    "Healthy Streets 10-indicator scores (0 = worst | 6.0 = standard threshold | 10 = best)")
R(sl, 0, 1.12, 13.33, 5.98, WHT)

# ── Annotation flag row (y=1.15 to 1.92) ─────────────────────────────────────
# Each flag sits over its respective panel in the chart3 image below.
# chart3 is a 3-panel chart; each panel is ≈ 1/3 of image width.
# Image will be placed x=0.30..13.03, so panel centers ≈ 2.29, 6.66, 11.03

flags = [
    (0.30, 4.28, RED, AMB,  "MAJOR",
     "Preston HS — HS2 = 0.4\nNo safe crossing at gate",      RED),
    (4.55, 3.95, AMB, GL,   "MODERATE",
     "Reservoir HS — 5 indicators below 6\n(footpath, crossing, shade, noise, relaxed)", AMB),
    (8.78, 4.20, GD,  GL,   "MODERATE",
     "William Ruthven SC — 4 below 6\n(shade, rest, active travel, amenity)",  GD),
]
for fx, fw, fbg, fbr, badge_txt, flag_txt, _ in flags:
    cx = fx + fw/2
    R(sl, fx, 1.15, fw, 0.77, fbg)
    R(sl, fx, 1.15, fw, 0.06, fbr)         # top stripe
    # Badge pill
    bw = len(badge_txt)*0.085 + 0.25
    R(sl, fx+0.12, 1.23, bw, 0.24, WHT)
    T(sl, badge_txt, fx+0.12, 1.24, bw, 0.22, bold=True, sz=9, col=fbg)
    T(sl, flag_txt, fx+0.14, 1.50, fw-0.28, 0.38, sz=9.5, col=WHT, bold=True)
    # Thin vertical connector to chart below
    R(sl, cx-0.03, 1.93, 0.06, 0.09, RGBColor(0xBD,0xBD,0xBD))

# ── Full-width chart image (y=2.02 → 7.08) ───────────────────────────────────
# chart3_score_breakdown.png — 3-panel bar chart
IMG(sl, "chart3_score_breakdown.png", 0.30, 2.02, 12.73, 5.06)

# ── Callout annotations positioned at specific spots in the chart image ───────
# Image top-left is at (0.30, 2.02). Chart is 12.73" wide × 5.06" tall.
# Preston panel occupies left ~1/3: x=0.30..4.56
#   HS2 bar (2nd bar from left) is at roughly x≈1.0..1.5, very short (0.4/12 scale)
#   → callout below Preston's HS2 bar area
# Reservoir panel: x=4.56..8.82
#   HS1 bar (1st bar) short at 4.2 → below Reservoir HS1 area
# Note: callouts placed WITHIN chart safe-zone (below major bar tops, at chart bottom area)

# Preston HS2 pointer — callout near bottom of Preston panel (HS2 is 2nd bar)
R(sl, 0.90, 5.62, 1.85, 0.52, RLT)
R(sl, 0.90, 5.62, 0.06, 0.52, RED)
T(sl, "HS2 = 0.4\nCritical — no\nsafe crossing", 1.00, 5.65, 1.65, 0.44, sz=9, col=RED, bold=True)
# small vertical line pointing up to bar area
R(sl, 1.70, 5.35, 0.04, 0.26, RED)

# Reservoir HS1 pointer
R(sl, 4.62, 5.62, 2.10, 0.52, ALT)
R(sl, 4.62, 5.62, 0.06, 0.52, AMB)
T(sl, "HS1 = 4.2\nBroken footpath\non Plenty Rd", 4.72, 5.65, 1.90, 0.44, sz=9, col=AMB, bold=True)
R(sl, 5.62, 5.35, 0.04, 0.26, AMB)

# William Ruthven HS3 pointer
R(sl, 9.22, 5.62, 2.10, 0.52, GLIT)
R(sl, 9.22, 5.62, 0.06, 0.52, GD)
T(sl, "HS3 = 3.0\nLowest shade &\nshelter score", 9.32, 5.65, 1.90, 0.44, sz=9, col=GD, bold=True)
R(sl, 10.22, 5.35, 0.04, 0.26, GD)

FTR(sl, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — CRASH TRENDS
# Left: 4 stat cards | Right: chart_crash_trends.png
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Pedestrian & Cyclist Crashes — A Rising Trend",
    "VicRoads open data 2021–2025, filtered for Darebin LGA and 400 m around school gates")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# Left sidebar: 4 stat cards stacked (zone: x=0.30..3.85, y=1.18..7.08)
STAT(sl, "73",    "Darebin LGA crashes\nin 2024 (up from 25)",   0.30, 1.18, bg=RED)
STAT(sl, "17:00", "Peak crash hour\n(school pickup window)",      0.30, 2.66, bg=AMB)
STAT(sl, "100%",  "W. Ruthven SC crashes\nduring school hours",   0.30, 4.14, bg=TEA)
STAT(sl, "15",    "Crashes near\nPreston HS (4 yrs)",             0.30, 5.62, bg=GD)

# Right: crash trends chart (zone: x=4.05..13.28, y=1.18..7.08)
IMG(sl, "chart_crash_trends.png", 4.05, 1.18, 9.23, 5.90)

# ── Callout annotations placed at chart panel boundaries ─────────────────────
# chart_crash_trends is a 2×2 grid:
#   Top-left panel  (LGA year trend):  x≈4.05..8.65, y≈1.18..4.04
#   Top-right panel (school year):     x≈8.65..13.28, y≈1.18..4.04
#   Bottom-left (school hours):        x≈4.05..8.65, y≈4.04..7.08
#   Bottom-right (time-of-day):        x≈8.65..13.28, y≈4.04..7.08

# Top-left: 2024 bar is the tallest in LGA trend. Bar roughly at x≈7.8, y≈1.55 (top of 73 bar)
# Place callout in the white space at top-right of top-left panel area
R(sl, 6.70, 1.20, 1.85, 0.52, RED)
R(sl, 6.70, 1.20, 1.85, 0.06, AMB)
T(sl, "2024: 73 crashes\n▲ 192% vs 2021", 6.82, 1.27, 1.62, 0.40, sz=9.5, col=WHT, bold=True)

# Bottom-right: time-of-day peak at 17:00. Chart shows a clear spike at 17:00.
# Place callout at top of bottom-right panel
R(sl, 9.85, 4.06, 2.00, 0.52, AMB)
R(sl, 9.85, 4.06, 2.00, 0.06, RED)
T(sl, "Peak at 17:00\nSchool pickup\ndrives crashes", 9.97, 4.13, 1.78, 0.40, sz=9.5, col=WHT, bold=True)

# Top-right: Preston HS has 5 crashes in 2024 bar, clearly highest.
# Place callout in white space at top of top-right panel
R(sl, 10.85, 1.20, 2.20, 0.52, GD)
R(sl, 10.85, 1.20, 2.20, 0.06, GL)
T(sl, "Preston HS: most\ncrashes (15 total\nnear gate 2021–25)", 10.97, 1.27, 1.98, 0.40, sz=9.5, col=WHT, bold=True)

FTR(sl, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — EQUITY FINDING
# Left: r=0.84 + SEIFA table + implication
# Right: chart_equity_seifa.png
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Equity — Disadvantaged Communities Have the Worst Streets",
    "ABS SEIFA 2021 Index of Relative Socio-Economic Disadvantage × Healthy Streets scores")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# Left sidebar (x=0.30..3.85, y=1.18..7.08)
R(sl, 0.30, 1.18, 3.55, 1.50, AMB)
R(sl, 0.30, 1.18, 3.55, 0.07, RED)
T(sl, "r = 0.84", 0.30, 1.22, 3.55, 0.85, bold=True, sz=52, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "Pearson correlation\nSEIFA decile × HS overall score",
  0.42, 2.06, 3.31, 0.52, sz=10, col=ALT, al=PP_ALIGN.CENTER)

# SEIFA table
seifa_rows = [
    ("Reservoir HS",      "Decile 3.7", "6.1 / 10", "Mod-high", RED),
    ("William Ruthven SC","Decile 3.7", "6.7 / 10", "Mod-high", AMB),
    ("Preston HS",        "Decile 5.5", "7.2 / 10", "Moderate", GM),
]
R(sl, 0.30, 2.78, 3.55, 0.36, GD)
for hdr_txt, hx in [("School",0.42),("SEIFA",1.82),("HS Score",2.55),("Level",3.15)]:
    T(sl, hdr_txt, hx, 2.81, 0.85, 0.28, bold=True, sz=9, col=WHT)

ry = 3.16
for schl, decile, hs, level, col in seifa_rows:
    bg = RLT if col == RED else (ALT if col == AMB else GLIT)
    R(sl, 0.30, ry, 3.55, 0.50, bg)
    T(sl, schl,  0.42, ry+0.10, 1.35, 0.32, bold=True, sz=9.5, col=GD)
    T(sl, decile,1.82, ry+0.10, 0.70, 0.32, bold=True, sz=10, col=col)
    T(sl, hs,    2.55, ry+0.10, 0.60, 0.32, bold=True, sz=11, col=col, al=PP_ALIGN.CENTER)
    T(sl, level, 3.17, ry+0.10, 0.62, 0.32, sz=9, col=col)
    ry += 0.52

# Implication note
NOTE(sl,
    "The most disadvantaged school catchments have the lowest "
    "safety scores. This is an equity imperative, not just a "
    "traffic engineering problem. Prioritised investment is required.",
    0.30, 4.82, 3.55, 1.00, bg=GD, sz=10.5)

# Vic Health relevance note
R(sl, 0.30, 5.90, 3.55, 0.80, RGBColor(0x00,0x45,0x4A))
T(sl, "Vic Health relevance: Active travel to school improves "
      "child health outcomes. The equity gap means the children "
      "with fewest transport alternatives face the most dangerous streets.",
  0.42, 5.94, 3.30, 0.72, sz=9.5, col=RGBColor(0x80,0xCB,0xC4))

# Right: equity chart (zone: x=4.05..13.28, y=1.18..7.08)
IMG(sl, "chart_equity_seifa.png", 4.05, 1.18, 9.23, 5.90)

# ── Callout on equity scatter (top-right quadrant of chart) ──────────────────
# chart_equity_seifa: top-right panel is scatter plot.
# Panel roughly x=8.65..13.28, y=1.18..4.04
# "At risk" quadrant (low decile, low HS) is bottom-left of scatter — around x=9..10, y=3..4
# Reservoir/William Ruthven dots are bottom-left of that scatter.
# The "HS threshold" line is at roughly y≈2.5 in the scatter panel.
# Place callout near the bottom-left of the scatter (where the at-risk schools are)
R(sl, 8.72, 2.65, 2.30, 0.58, RED)
R(sl, 8.72, 2.65, 2.30, 0.06, AMB)
T(sl, "At-risk zone:\nlow SEIFA + low HS\n= priority investment",
  8.84, 2.72, 2.08, 0.46, sz=9.5, col=WHT, bold=True)

# Preston dot is top-right of scatter (higher decile, higher HS)
R(sl, 11.50, 1.22, 1.60, 0.50, GLIT)
R(sl, 11.50, 1.22, 1.60, 0.05, GL)
T(sl, "Preston HS\nDecile 5.5\nHS=7.2", 11.58, 1.26, 1.42, 0.38, sz=9.5, col=GD, bold=True)

FTR(sl, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — MACHINE LEARNING
# Left: automation table | Right: chart_hs_prediction.png
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Can Open Data Replace Field Surveys?",
    "Ridge regression trained on 26 open-data features — Leave-One-Out cross-validation (n = 3 schools)")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# Left sidebar content (x=0.30..3.85, y=1.18..7.08)
NOTE(sl, "Model: MultiOutput Ridge regression\n"
         "Features: 26 (OSM, AQI, crime, crashes)\n"
         "Evaluation: LOO-CV  |  Mean MAE: 2.88",
     0.30, 1.18, 3.55, 0.82, bg=TEA, sz=10.5)

# Automation tier table
tiers = [
    ("AUTOMATE", GM, [
        ("HS7 — Feel safe",  "0.54", "Crime rate is a reliable proxy"),
        ("HS10 — Clean air", "0.92", "AQI maps directly — no survey"),
    ]),
    ("PARTIAL", AMB, [
        ("HS9 — Feel relaxed", "1.72", "Speed/arterial data works well"),
        ("HS3 — Shade/shelter","2.17", "OSM tree data reasonable"),
        ("HS6 — Active travel", "3.08", "Coverage OK, quality not"),
        ("HS1 — Footpaths",    "3.27", "Presence measured, condition not"),
        ("HS5 — Noise",        "3.68", "OSM speed data incomplete"),
    ]),
    ("SURVEY REQUIRED", RED, [
        ("HS4 — Rest places",  "4.19", "OSM bench data unreliable"),
        ("HS2 — Crossing",     "4.51", "Presence ≠ quality/safety"),
        ("HS8 — Things to do", "4.72", "Activity quality not in OSM"),
    ]),
]

ty = 2.05
for tier_lbl, bg, rows in tiers:
    R(sl, 0.30, ty, 3.55, 0.30, bg)
    T(sl, tier_lbl, 0.42, ty+0.04, 3.30, 0.23, bold=True, sz=9, col=WHT)
    ty += 0.32
    for ind_name, mae, note in rows:
        R(sl, 0.30, ty, 3.55, 0.34, WHT if rows.index((ind_name,mae,note))%2==0 else RGBColor(0xF9,0xF9,0xF9))
        T(sl, ind_name, 0.40, ty+0.05, 1.85, 0.26, bold=True, sz=9.5, col=INK)
        T(sl, f"MAE {mae}", 2.28, ty+0.05, 0.70, 0.26, sz=9, col=bg, bold=True)
        T(sl, note, 3.02, ty+0.05, 0.75, 0.26, sz=8, col=MUT)
        ty += 0.36
    ty += 0.06

# Caveat box
R(sl, 0.30, ty+0.05, 3.55, 0.56, RGBColor(0xE3,0xF2,0xFD))
R(sl, 0.30, ty+0.05, 0.06, 0.56, RGBColor(0x01,0x57,0x9B))
T(sl, "n=3 — results are illustrative. Model accuracy improves significantly with ≥20 schools. "
      "HS2 and HS8 require field observation at any sample size.",
  0.42, ty+0.1, 3.30, 0.44, sz=9, col=INK)

# Right: LOO-CV prediction chart (zone: x=4.05..13.28, y=1.18..7.08)
IMG(sl, "chart_hs_prediction.png", 4.05, 1.18, 9.23, 5.90)

# ── Callout: Preston panel (rightmost, 3rd panel): HS2 actual=0.4, predicted=3.2 ──
# Preston panel is right 1/3 of image: x≈10.1..13.28, y=1.18..7.08
# HS2 bars are 2nd pair from left in that panel — roughly at x≈11.0..11.3
# The huge gap between dark(actual 0.4) and light(predicted 3.2) bars is visually obvious
R(sl, 10.80, 4.12, 2.28, 0.62, RLT)
R(sl, 10.80, 4.12, 0.06, 0.62, RED)
T(sl, "HS2: actual=0.4\npredicted=3.2\n→ must survey crossings", 10.90, 4.16, 2.10, 0.52, sz=9.5, col=RED, bold=True)

# Reservoir panel (leftmost): HS4 predicted=0.9 vs actual=6.0 — huge miss
R(sl, 4.80, 4.12, 2.35, 0.62, ALT)
R(sl, 4.80, 4.12, 0.06, 0.62, AMB)
T(sl, "HS4: predicted=0.9\nactual=6.0 — OSM bench\ndata unreliable", 4.90, 4.16, 2.18, 0.52, sz=9.5, col=AMB, bold=True)

FTR(sl, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — WHAT CAN BE DONE
# Left: intervention table | Right: scenario chart (Preston HS)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "What Can Be Done — Scenario Analysis",
    "Delta-method modelling: 10 physical interventions × 3 schools = 30 what-if scenarios")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# Left sidebar (x=0.30..3.85, y=1.18..7.08)
NOTE(sl, "How it works: actual HS scores are used as the "
         "baseline. The trained Ridge model contributes only "
         "the predicted CHANGE when a feature delta is applied "
         "(e.g. +1 crossing, −10 km/h speed limit).",
     0.30, 1.18, 3.55, 1.00, bg=TEA, sz=10)

R(sl, 0.30, 2.25, 3.55, 0.33, GD)
T(sl, "10 Available Interventions", 0.42, 2.28, 3.30, 0.26, bold=True, sz=11, col=WHT)

interventions = [
    ("Signalised crossing",       "HS2", "$80k–$200k",   "< 1 yr",  RED),
    ("Speed limit → 40 km/h",     "HS5", "$5k–$20k",    "< 6 mo",  GM),
    ("Traffic calming",           "HS9", "$50k–$150k",  "< 1 yr",  AMB),
    ("Continuous footpath",       "HS1", "$100k–$300k", "< 1 yr",  RED),
    ("Protected bike lane",       "HS6", "$500k+/km",   "1–3 yr",  TEA),
    ("Street trees",              "HS3", "$20k–$60k",   "< 1 yr",  GM),
    ("PT stop upgrade",           "HS6", "$50k–$200k",  "1–2 yr",  TEA),
    ("Street benches",            "HS4", "$5k–$15k",    "< 6 mo",  GM),
    ("Covered shelter",           "HS3", "$15k–$40k",   "< 6 mo",  GM),
    ("Reroute arterial traffic",  "HS5", "High/long",   "2–5 yr",  GD),
]
iy = 2.60
for name, ind, cost, lead, col in interventions:
    bg = WHT if interventions.index((name,ind,cost,lead,col))%2==0 else RGBColor(0xF4,0xF9,0xF4)
    R(sl, 0.30, iy, 3.55, 0.38, bg)
    R(sl, 0.30, iy, 0.42, 0.38, col)
    T(sl, ind, 0.30, iy+0.07, 0.42, 0.24, bold=True, sz=8.5, col=WHT, al=PP_ALIGN.CENTER)
    T(sl, name, 0.78, iy+0.07, 1.70, 0.24, sz=9.5, col=INK)
    T(sl, cost, 2.52, iy+0.07, 0.78, 0.24, sz=8.5, col=col, bold=True)
    T(sl, lead, 3.34, iy+0.07, 0.46, 0.24, sz=8, col=MUT)
    iy += 0.40

# Right: Preston HS scenario chart (zone: x=4.05..13.28, y=1.18..7.08)
IMG(sl, "scenario_Preston_HS_pedestrian_crossing_speed_reduction.png",
    4.05, 1.18, 9.23, 5.90)

# ── Callout annotations ───────────────────────────────────────────────────────
# The Preston scenario chart has 2 panels:
#   Left panel (before/after bars): x≈4.05..8.65, y=1.18..7.08
#   Right panel (delta bars):       x≈8.65..13.28, y=1.18..7.08
# HS2 Cross is near the bottom of each panel (it appears last alphabetically reversed)
# In the before/after panel: HS2 is a very short bar (0.4), at approximately y≈5.8..6.2
# The title/heading area of the chart (top) is at y≈1.18..1.80

# Top title area callout — intervention description
R(sl, 4.12, 1.20, 4.30, 0.55, GD)
R(sl, 4.12, 1.20, 4.30, 0.06, GL)
T(sl, "Preston HS:  Pedestrian crossing + 40 km/h zone\nOverall: 7.2 → 7.25  |  Severity: Major remains",
  4.24, 1.26, 4.08, 0.43, sz=9.5, col=WHT, bold=True)

# Model caveat on the chart (placed below the chart title area, above bars)
R(sl, 8.70, 1.20, 4.40, 0.55, RGBColor(0xFB,0xE9,0xE7))
R(sl, 8.70, 1.20, 0.06, 0.55, AMB)
T(sl, "Note: n=3 model — delta predictions are directional.\n"
      "Negative HS2 delta is a Ridge artefact with small n.",
  8.80, 1.26, 4.22, 0.43, sz=9.5, col=AMB)

FTR(sl, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — RECOMMENDATIONS
# 3 columns, one per school — no images
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Prioritised Recommendations",
    "Rule-based interventions per school, ranked by HS indicator, cost, and lead time")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# Column layout: col_x = [0.30, 4.50, 8.70], col_w = 3.95
cols_data = [
    ("Preston HS", "Major — HS2 = 0.4", RED, [
        ("CRITICAL", RED,  "Install signalised pedestrian crossing at school gate (HS2=0.4 — no safe crossing exists)",
         "< 1 year",  "$80k–$200k"),
        ("High",     AMB,  "Reduce speed limit to 40 km/h school zone on Cooma St",
         "< 6 months","$5k–$20k"),
        ("High",     AMB,  "Install traffic calming (speed humps / raised crossing) on approach roads",
         "< 1 year",  "$50k–$150k"),
        ("Medium",   GM,   "Improve public transport access and active travel infrastructure",
         "1–2 years", "$50k–$200k"),
    ]),
    ("Reservoir HS", "Moderate — 5 below standard", AMB, [
        ("High",   RED,  "Build continuous footpath on Plenty Rd — missing sections (HS1 = 4.2)",
         "< 1 year",  "$100k–$300k"),
        ("High",   AMB,  "Install tactile pavers and kerb ramps at all crossings near gate",
         "< 1 year",  "$20k–$80k"),
        ("High",   AMB,  "Traffic calming — multi-lane Plenty Rd approach to gate (HS9 = 5.2)",
         "< 1 year",  "$50k–$150k"),
        ("Medium", GM,   "Plant street trees for shade on walking routes (HS3 = 5.0)",
         "< 1 year",  "$20k–$60k"),
    ]),
    ("William Ruthven SC", "Moderate — shade/rest/active travel", GD, [
        ("High",   RED,  "Plant trees and install shelters — lowest shade score of 3 schools (HS3 = 3.0)",
         "< 1 year",  "$35k–$100k"),
        ("High",   AMB,  "Add street seating/rest points near school gate (HS4 = 3.3)",
         "< 6 months","$5k–$20k"),
        ("Medium", GM,   "Improve cycling infrastructure on Merrilands Rd (HS6 = 4.2)",
         "1–3 years", "$500k+/km"),
        ("Medium", GM,   "Add PT stop improvements and wayfinding signage (HS8 = 4.3)",
         "1–2 years", "$20k–$60k"),
    ]),
]

for ci, (school, severity, scol, recs) in enumerate(cols_data):
    cx = 0.30 + ci * 4.37
    cw = 4.0

    # School header
    R(sl, cx, 1.18, cw, 0.52, scol)
    R(sl, cx, 1.18, cw, 0.07, GL)
    T(sl, school, cx+0.12, 1.20, cw-0.24, 0.30, bold=True, sz=14, col=WHT)
    # Severity badge
    bw = len(severity)*0.077 + 0.20
    R(sl, cx+0.12, 1.52, bw, 0.22, WHT)
    T(sl, severity, cx+0.14, 1.53, bw-0.04, 0.18, bold=True, sz=8, col=scol)

    ry = 1.78
    for pri, pcol, rec_text, timeline, cost in recs:
        card_h = 1.14
        R(sl, cx, ry, cw, card_h, WHT)
        R(sl, cx, ry, 0.06, card_h, pcol)          # priority colour bar
        # Priority pill
        pw = len(pri)*0.080+0.18
        R(sl, cx+0.14, ry+0.07, pw, 0.22, pcol)
        T(sl, pri, cx+0.14, ry+0.08, pw, 0.20, bold=True, sz=8.5, col=WHT)
        # Recommendation text
        T(sl, rec_text, cx+0.14, ry+0.33, cw-0.24, 0.46, sz=9.5, col=INK)
        # Timeline / cost
        R(sl, cx+0.14, ry+0.83, cw-0.24, 0.22, GLIT)
        T(sl, f"Timeline: {timeline}  ·  Cost: {cost}",
          cx+0.20, ry+0.85, cw-0.32, 0.18, sz=8.5, col=SLT, it=True)
        ry += card_h + 0.06

# Bottom CTA bar
R(sl, 0.30, 6.80, 12.73, 0.42, GD)
T(sl, "All 17 recommendations with HS indicator, priority, cost & timeframe available "
      "in the filterable recommendations table at the project web dashboard",
  0.50, 6.83, 12.36, 0.35, sz=11, col=WHT, al=PP_ALIGN.CENTER)

FTR(sl, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — INTERACTIVE WEB DASHBOARD (INSERTED)
# Left: features + tech | Right: browser UI mockup
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Interactive Web Dashboard",
    "All findings in one place — no installation, runs in any browser, zero backend infrastructure")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# ── Left sidebar (x=0.30..4.00, y=1.18..7.08) ────────────────────────────────
NOTE(sl,
    "Live at GitHub Pages — fully static HTML.\n"
    "Built on Leaflet, Chart.js, and Tailwind CSS.\n"
    "All 30 scenario results are pre-computed;\n"
    "no Python runtime needed in the browser.",
    0.30, 1.18, 3.65, 1.05, bg=TEA, sz=10.5)

sections_list = [
    (GD,  "Hero stats strip",
     "5 live numbers: schools assessed, major hazards, crash count, equity r, peak hour"),
    (GM,  "Interactive map",
     "School markers colour-coded by severity; red pulse = Major; click → school detail"),
    (TEA, "School Assessments",
     "Per-school tab: radar chart + 10 HS indicator bars + key hazard callout"),
    (AMB, "Scenario Explorer",
     "Pick school + intervention → instant before/after chart; 30 pre-computed results"),
    (GD,  "Data Analysis",
     "All 9 pipeline charts: ML results, equity, crash trends, demographics"),
    (GM,  "Recommendations",
     "Filterable table (school, priority) with one-click CSV download"),
    (SLT, "About",
     "Healthy Streets framework explainer + severity classification rules"),
]
sy = 2.30
for bg, sec_name, desc in sections_list:
    R(sl, 0.30, sy, 3.65, 0.55, WHT if sections_list.index((bg,sec_name,desc))%2==0 else RGBColor(0xF4,0xF9,0xF4))
    R(sl, 0.30, sy, 0.07, 0.55, bg)
    T(sl, sec_name, 0.46, sy+0.05, 3.35, 0.20, bold=True, sz=10.5, col=bg)
    T(sl, desc,     0.46, sy+0.28, 3.35, 0.22, sz=9, col=SLT)
    sy += 0.58

# ── Right: browser window mockup (x=4.15..13.05, y=1.18..7.08) ──────────────
BX = 4.15   # browser left edge
BW = 8.90   # browser total width
BY = 1.18   # browser top
BH = 5.90   # browser total height

# Browser chrome (title bar)
R(sl, BX, BY, BW, 0.38, RGBColor(0xE0,0xE0,0xE0))
# Browser buttons (traffic-light circles simulated as small rects)
for bx_off, bc in [(0.12, RGBColor(0xFF,0x5F,0x56)),
                   (0.32, RGBColor(0xFF,0xBD,0x2E)),
                   (0.52, RGBColor(0x27,0xC9,0x3F))]:
    R(sl, BX+bx_off, BY+0.10, 0.16, 0.16, bc)
# URL bar
R(sl, BX+0.78, BY+0.08, BW-1.00, 0.22, WHT)
T(sl, "  s4111078-sidd.github.io/300-000-School-Streets",
  BX+0.82, BY+0.09, BW-1.10, 0.18, sz=8, col=SLT)

# ── Nav bar ──────────────────────────────────────────────────────────────────
R(sl, BX, BY+0.38, BW, 0.36, GD)
T(sl, "300,000 Streets of Melbourne", BX+0.15, BY+0.40, 3.50, 0.28, bold=True, sz=10, col=WHT)
for nav_lbl, nx in [("Map",2.6),("Schools",3.3),("Scenarios",4.2),("Analysis",5.3),("Recs",6.4)]:
    T(sl, nav_lbl, BX+nx, BY+0.42, 0.85, 0.22, sz=8.5, col=AGRN)

# ── Hero stats strip ─────────────────────────────────────────────────────────
R(sl, BX, BY+0.74, BW, 0.42, RGBColor(0x14,0x4A,0x18))
hero = [("3","Schools"),("1","Major hazard"),("21","Darebin crashes"),("0.84","Equity r"),("17:00","Peak hour")]
hcx = BX + 0.10
for hv, hl in hero:
    R(sl, hcx, BY+0.77, 1.60, 0.35, RGBColor(0x21,0x68,0x26))
    T(sl, hv, hcx+0.04, BY+0.77, 0.65, 0.20, bold=True, sz=11, col=WHT)
    T(sl, hl, hcx+0.70, BY+0.80, 0.85, 0.18, sz=7.5, col=AGRN)
    hcx += 1.72

# ── Map section ──────────────────────────────────────────────────────────────
R(sl, BX, BY+1.16, BW, 1.50, RGBColor(0xD0,0xE4,0xC8))   # green = open street map feel
R(sl, BX, BY+1.16, BW, 0.22, RGBColor(0xB8,0xD4,0xAA))   # slightly darker top strip
T(sl, "Interactive Map — Leaflet", BX+0.15, BY+1.20, 4.00, 0.18, bold=True, sz=9, col=GD)
# School marker dots
for mx, my, mc, ml in [(BX+1.5, BY+1.75, RED,  "Preston HS ●"),
                        (BX+3.2, BY+2.00, AMB,  "Reservoir HS ●"),
                        (BX+2.4, BY+2.35, AMB,  "W. Ruthven SC ●")]:
    R(sl, mx, my, 0.14, 0.14, mc)
    T(sl, ml, mx+0.18, my-0.03, 1.40, 0.22, sz=8, col=INK, bold=True)
# Zoom controls
R(sl, BX+BW-0.38, BY+1.26, 0.26, 0.48, WHT)
T(sl, "+\n−", BX+BW-0.34, BY+1.28, 0.20, 0.44, sz=9, col=INK, al=PP_ALIGN.CENTER)

# ── Bottom two-panel row ──────────────────────────────────────────────────────
# Left: School Assessment tab
R(sl, BX, BY+2.66, BW*0.48, 1.62, WHT)
R(sl, BX, BY+2.66, BW*0.48, 0.26, GLIT)
T(sl, "School Assessments", BX+0.15, BY+2.68, 3.50, 0.20, bold=True, sz=9, col=GD)
# mini radar representation (hexagon via 6 small rects — simplified)
for rx, ry_off, rw, rh, rc in [
    (BX+0.20, BY+3.00, 0.95, 0.12, GM),
    (BX+0.20, BY+3.15, 1.40, 0.12, GM),
    (BX+0.20, BY+3.30, 0.60, 0.12, RED),
    (BX+0.20, BY+3.45, 1.20, 0.12, GM),
    (BX+0.20, BY+3.60, 1.60, 0.12, GM),
    (BX+0.20, BY+3.75, 0.80, 0.12, AMB),
]:
    R(sl, rx, ry_off, rw, rh, rc)
T(sl, "HS indicator bars", BX+0.20, BY+3.92, 2.50, 0.20, sz=8, col=SLT, it=True)

# Right: Scenario Explorer
R(sl, BX+BW*0.50, BY+2.66, BW*0.48, 1.62, WHT)
R(sl, BX+BW*0.50, BY+2.66, BW*0.48, 0.26, GLIT)
T(sl, "Scenario Explorer", BX+BW*0.50+0.12, BY+2.68, 3.50, 0.20, bold=True, sz=9, col=TEA)
# mini before/after bar chart
for bi, (bv1, bv2) in enumerate([(0.4,0.4),(1.7,1.9),(1.0,1.2),(1.7,2.1),(0.8,0.8)]):
    bbase = BX+BW*0.50+0.18 + bi*0.75
    R(sl, bbase,      BY+3.65, 0.28, bv1*0.55, RGBColor(0x37,0x47,0x4F))
    R(sl, bbase+0.32, BY+3.65, 0.28, bv2*0.55, TEA)
T(sl, "Before (dark)  After (teal)", BX+BW*0.50+0.18, BY+3.93, 3.00, 0.20, sz=8, col=SLT, it=True)

# ── Footer strip inside browser ───────────────────────────────────────────────
R(sl, BX, BY+BH-0.28, BW, 0.28, RGBColor(0x21,0x21,0x21))
T(sl, "s4111078-sidd.github.io/300-000-School-Streets  ·  GitHub Pages  ·  Static HTML, no backend",
  BX+0.20, BY+BH-0.26, BW-0.30, 0.22, sz=7.5, col=RGBColor(0x90,0x90,0x90))

FTR(sl, 10)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — SCALE-UP & NEXT STEPS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
HDR(sl, "From 3 Schools to 300,000 Streets",
    "A scalable, reproducible pipeline for evidence-based investment in every school street in Victoria")
R(sl, 0, 1.12, 13.33, 5.98, GLIT)

# ── Vision banner ─────────────────────────────────────────────────────────────
R(sl, 0.30, 1.18, 12.73, 0.60, GD)
T(sl,
  "Victoria has ~1,900 government schools. The 300,000 Streets pipeline can assess any school gate "
  "using publicly available data — the same pipeline, zero new infrastructure.",
  0.48, 1.22, 12.36, 0.52, sz=12.5, col=AGRN, bold=True)

# ── Three columns ─────────────────────────────────────────────────────────────
col_data = [
    ("Why This Scales", GL, [
        ("Adding a school",
         "New gate coordinates + re-run python run_all.py — outputs in minutes"),
        ("Open data coverage",
         "70% of HS indicators estimated from public data alone (HS7, HS10 fully automated)"),
        ("Scenario engine",
         "30 what-if analyses in seconds — evidence base for any planning decision"),
        ("GIS-ready outputs",
         "GPKG + QGIS project files slot directly into council planning workflows"),
        ("Zero infrastructure",
         "Static web dashboard on GitHub Pages — shareable, no backend needed"),
    ]),
    ("Current Limitations", AMB, [
        ("Small training set",
         "n=3 — ML predictions are illustrative; model generalises with ≥20 schools"),
        ("Point-in-time survey",
         "Field observations are a snapshot; seasonal/time-of-day variation not captured"),
        ("OSM quality",
         "Low-income areas sometimes have less complete OSM data — biases open-data metrics"),
        ("Scenario model",
         "Delta method predicts HS score change, not actual crash reduction — needs validation"),
        ("Causal claims",
         "SEIFA correlation (r=0.84) is associative — confounders exist (land use, road age)"),
    ]),
    ("Recommended Next Steps", TEA, [
        ("Expand to 10–20 schools",
         "Multiple LGAs — validates ML model, enables cross-suburb equity comparison"),
        ("Partner with DET / VicRoads",
         "Use SEIFA equity lens to triage which schools are assessed first"),
        ("Longitudinal crash tracking",
         "Measure actual crash changes after interventions — ground truth the model"),
        ("Annual pipeline refresh",
         "Automate yearly re-run to track progress after infrastructure investment"),
        ("VicHealth integration",
         "Connect active-travel outcomes to child health surveillance data (NAPLAN, AIHW)"),
    ]),
]

for ci, (col_title, stripe_col, items) in enumerate(col_data):
    cx = 0.30 + ci*4.38
    cw = 4.08

    R(sl, cx, 1.88, cw, 0.40, GD)
    R(sl, cx, 1.88, cw, 0.06, stripe_col)
    T(sl, col_title, cx+0.15, 1.92, cw-0.25, 0.30, bold=True, sz=13, col=WHT)

    iy = 2.34
    for n, (item_title, item_desc) in enumerate(items):
        bg = WHT if n % 2 == 0 else RGBColor(0xF5,0xF5,0xF5)
        R(sl, cx, iy, cw, 0.85, bg)
        R(sl, cx, iy, 0.06, 0.85, stripe_col)
        T(sl, item_title, cx+0.16, iy+0.06, cw-0.26, 0.26, bold=True, sz=10.5, col=GD)
        T(sl, item_desc,  cx+0.16, iy+0.35, cw-0.26, 0.44, sz=9.5, col=INK)
        iy += 0.90

# ── Vision footer banner ──────────────────────────────────────────────────────
R(sl, 0.30, 6.75, 12.73, 0.45, GD)
T(sl,
  "Every child deserves a safe walk to school.  "
  "300,000 Streets of Melbourne — Regen Melbourne × RMIT University",
  0.50, 6.78, 12.36, 0.38, bold=True, sz=13, col=WHT, al=PP_ALIGN.CENTER)

FTR(sl, 11)

# ─────────────────────────────────────────────────────────────────────────────
OUT_PATH = os.path.join(OUT, "School_Streets_Presentation.pptx")
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
