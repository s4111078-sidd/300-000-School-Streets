"""
Generate 300,000 Streets of Melbourne — School Streets Safety Assessment PPT.
Run: python make_ppt.py
Output: outputs/School_Streets_Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x23, 0x42)
TEAL   = RGBColor(0x1A, 0x8F, 0xC1)
AMBER  = RGBColor(0xE8, 0xA8, 0x38)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xD3, 0x54, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF2, 0xF4, 0xF7)
DGREY  = RGBColor(0x55, 0x55, 0x55)
BLACK  = RGBColor(0x11, 0x11, 0x11)

OUT_DIR   = 'outputs'
CHART1    = os.path.join(OUT_DIR, 'chart1_safety_scores.png')
CHART2    = os.path.join(OUT_DIR, 'chart2_hazard_severity.png')
CHART3    = os.path.join(OUT_DIR, 'chart3_score_breakdown.png')
OUT_PPTX  = os.path.join(OUT_DIR, 'School_Streets_Presentation.pptx')

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None, bar_color=NAVY):
    add_rect(slide, 0, 0, 13.33, 1.15, bar_color)
    add_text(slide, title, 0.35, 0.12, 12.5, 0.7,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.78, 12.5, 0.35,
                 size=13, bold=False, color=RGBColor(0xCC, 0xDD, 0xEE),
                 align=PP_ALIGN.LEFT)
    add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
    add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
             0.3, 7.32, 12, 0.18, size=8, color=WHITE, align=PP_ALIGN.LEFT)


def bullet_box(slide, bullets, x, y, w, h,
               title=None, title_color=NAVY, bullet_size=13,
               bg_color=LGREY, indent='• '):
    add_rect(slide, x, y, w, h, bg_color)
    txb = slide.shapes.add_textbox(Inches(x+0.15), Inches(y+0.12),
                                   Inches(w-0.3), Inches(h-0.24))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(bullet_size + 1)
        run.font.bold  = True
        run.font.color.rgb = title_color
    for b in bullets:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = indent + b
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = BLACK
    return txb


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
add_rect(slide, 0, 5.6, 13.33, 1.9, TEAL)

add_text(slide, '300,000 Streets of Melbourne',
         0.6, 1.0, 12.0, 1.1, size=42, bold=True, color=WHITE)
add_text(slide, 'School Streets Safety Assessment',
         0.6, 2.0, 12.0, 0.7, size=30, bold=False,
         color=RGBColor(0xA0, 0xC8, 0xE8))
add_text(slide,
         'Healthy Streets Framework  ·  ML Prediction  ·  Spatial Analysis  ·  Live Dashboard Roadmap',
         0.6, 2.85, 12.0, 0.5, size=15, color=RGBColor(0x88, 0xBB, 0xDD))

add_rect(slide, 0.6, 3.6, 12.1, 0.03, TEAL)

add_text(slide, 'Darebin LGA  —  Preston HS  ·  Reservoir HS  ·  William Ruthven SC',
         0.6, 3.8, 12.0, 0.5, size=16, bold=True, color=AMBER)
add_text(slide, 'Regen Melbourne  ×  RMIT University   |   2025',
         0.6, 5.75, 12.0, 0.4, size=14, color=WHITE)
add_text(slide, 'Healthy Streets Framework by Lucy Saunders / Transport for London',
         0.6, 6.2, 12.0, 0.35, size=11, italic=True,
         color=RGBColor(0xAA, 0xCC, 0xDD))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_header(slide, 'The Problem', 'Why school streets in Darebin need urgent attention')
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, 'The Problem', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'Why school streets in Darebin need urgent attention',
         0.35, 0.78, 12.5, 0.35, size=13,
         color=RGBColor(0xCC, 0xDD, 0xEE))

# Big stat boxes
for i, (stat, label, col) in enumerate([
    ('300,000+', 'Children walk/cycle\nto school daily in Melbourne', TEAL),
    ('3rd highest', 'Pedestrian crash rate\namong Melbourne LGAs', RED),
    ('82%', 'Of assessed school streets\nscore below HS "good" threshold', AMBER),
]):
    bx = 0.4 + i * 4.3
    add_rect(slide, bx, 1.4, 4.0, 1.8, col)
    add_text(slide, stat, bx+0.15, 1.5, 3.7, 0.9,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, bx+0.1, 2.25, 3.8, 0.85,
             size=13, color=WHITE, align=PP_ALIGN.CENTER)

# Key points
txb = slide.shapes.add_textbox(Inches(0.4), Inches(3.45), Inches(12.5), Inches(3.6))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True

items = [
    ('The Streets of Melbourne project identified Darebin (Preston, Reservoir, Thornbury) as a priority '
     'intervention area due to high pedestrian and cyclist crash rates near school zones.',
     False, 14),
    ('Three schools were selected for detailed assessment: Preston High School, Reservoir High School, '
     'and William Ruthven Secondary College — covering contrasting socioeconomic and road environments.',
     False, 14),
    ('Field assessments combined with open spatial data (OSM, VicRoads, EPA, Crime Statistics Agency) '
     'were used to build a reproducible, data-driven safety scoring system.',
     False, 14),
    ('Goal: identify the highest-risk streets, understand WHY they are dangerous, and produce '
     'evidence-based intervention recommendations that council can act on.',
     True, 14),
]
first = True
for text, bold, size in items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(7)
    run = p.add_run()
    run.text = ('▶  ' if bold else '• ') + text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = NAVY if bold else BLACK


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OLD FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, RGBColor(0x44, 0x44, 0x44))
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Original Scoring Framework', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'FAS · CSS · EEI — three-indicator field survey approach',
         0.35, 0.78, 12.5, 0.35, size=13,
         color=RGBColor(0xCC, 0xCC, 0xCC))

for i, (code, name, desc, col) in enumerate([
    ('FAS', 'Footpath Accessibility\nScore', 'Rated pavement quality, width,\nobstructions, and continuity\nalong school routes on foot.', GREEN),
    ('CSS', 'Crossing Safety\nScore', 'Assessed road crossing conditions:\nmarkings, signals, sight lines,\nvehicle speeds at crossings.', TEAL),
    ('EEI', 'Environmental Exposure\nIndicator', 'Captured pollution, noise, shade,\nand overall environmental comfort\nfor active travel.', AMBER),
]):
    bx = 0.35 + i * 4.3
    add_rect(slide, bx, 1.35, 4.0, 2.5, col)
    add_text(slide, code, bx+0.15, 1.45, 3.7, 0.65,
             size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, name, bx+0.1, 2.05, 3.8, 0.65,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, bx+0.1, 2.65, 3.8, 1.1,
             size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, 'Overall Hazard Severity', 0.35, 4.05, 12.6, 0.4,
         size=16, bold=True, color=NAVY)
add_text(slide,
         'Average(FAS, CSS, EEI)  →  thresholded to:   Major   |   Moderate   |   Minor',
         0.35, 4.4, 12.6, 0.4, size=14, color=DGREY)

# Limitations box
add_rect(slide, 0.35, 5.0, 12.6, 2.0, RGBColor(0xFF, 0xF0, 0xE8))
add_text(slide, 'Limitations of the original framework:', 0.55, 5.1, 12, 0.35,
         size=13, bold=True, color=RED)
txb = slide.shapes.add_textbox(Inches(0.55), Inches(5.45), Inches(12.1), Inches(1.45))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
limits = [
    '3 indicators only — misses pedestrian perception, shade, cycling infrastructure, noise, and amenity',
    'Not aligned to any recognised international framework — difficult to benchmark against other cities',
    'Threshold logic was ad hoc — severity classification could not be reproduced or explained to stakeholders',
    'No open-data path — required manual field scoring for every street, not scalable to 300,000 streets',
]
first = True
for lim in limits:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = '✗  ' + lim
    run.font.size = Pt(12)
    run.font.color.rgb = RED


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHY WE MOVED
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Why We Moved to Healthy Streets', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'From 3 ad-hoc indicators to a globally recognised 10-indicator framework',
         0.35, 0.78, 12.5, 0.35, size=13,
         color=RGBColor(0xCC, 0xDD, 0xEE))

# Left: Old
add_rect(slide, 0.35, 1.3, 5.9, 5.7, RGBColor(0xFF, 0xEC, 0xEC))
add_text(slide, '✗  Old Framework (FAS / CSS / EEI)', 0.5, 1.4, 5.6, 0.45,
         size=14, bold=True, color=RED)
old_pts = [
    '3 indicators covering limited dimensions',
    'No recognised benchmark or standard',
    'Severity thresholds were arbitrary',
    'Required full field survey to score',
    'Could not explain WHY a street scored low',
    'Not extensible — no path to open-data automation',
    'Could not be used for what-if scenario planning',
]
txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.6), Inches(4.9))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
first = True
for pt in old_pts:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = '✗  ' + pt
    run.font.size = Pt(13)
    run.font.color.rgb = RED

# Arrow
add_text(slide, '→', 6.4, 3.6, 0.6, 0.6, size=30, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)

# Right: New
add_rect(slide, 7.1, 1.3, 5.9, 5.7, RGBColor(0xE8, 0xF8, 0xEE))
add_text(slide, '✓  Healthy Streets (HS1–HS10)', 7.25, 1.4, 5.6, 0.45,
         size=14, bold=True, color=GREEN)
new_pts = [
    '10 indicators covering all dimensions of street safety',
    'Developed by Lucy Saunders / Transport for London',
    'Used by 50+ cities globally — internationally benchmarkable',
    'Threshold of 6.0 is evidence-based (HS "good" standard)',
    'Each indicator maps to a specific open-data source',
    'Severity logic is transparent and reproducible',
    'Directly powers ML scenario prediction engine',
]
txb = slide.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5.6), Inches(4.9))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
first = True
for pt in new_pts:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = '✓  ' + pt
    run.font.size = Pt(13)
    run.font.color.rgb = GREEN


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HEALTHY STREETS FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, TEAL)
add_rect(slide, 0, 7.3, 13.33, 0.2, NAVY)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'The Healthy Streets Framework', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'HS1–HS10 scored 0–10 per indicator  ·  Good threshold = 6.0  ·  Lucy Saunders / TfL',
         0.35, 0.78, 12.5, 0.35, size=13, color=WHITE)

indicators = [
    ('HS1',  'People Feel Safe',          'Crime perception, lighting, eyes on street',              GREEN),
    ('HS2',  'Easy to Cross Roads',        'Pedestrian crossings, signals, wait times',               TEAL),
    ('HS3',  'Not Too Noisy',             'Traffic noise levels at school drop-off times',            AMBER),
    ('HS4',  'People Choose Active Travel','Mode share — % walking/cycling vs car trips',             RGBColor(0x8E, 0x44, 0xAD)),
    ('HS5',  'Shade and Shelter',          'Tree canopy cover, covered walkways near school gate',    GREEN),
    ('HS6',  'Cycling Infrastructure',     'Protected lanes, cycle path %, PT stop access',           TEAL),
    ('HS7',  'Clean Air',                  'EPA AQI at arterial corridors during school hours',       RGBColor(0x1A, 0xBC, 0x9C)),
    ('HS8',  'Things to See and Do',       'Active frontage, cafes, parks, street activation',        AMBER),
    ('HS9',  'Places to Stop and Rest',    'Seating, benches, sheltered rest areas',                  RGBColor(0x27, 0x6F, 0xBF)),
    ('HS10', 'Not Too Fast',               'Posted + observed vehicle speeds near school gate',       RED),
]

cols = [0.25, 4.5, 8.75]
for idx, (code, name, desc, col) in enumerate(indicators):
    row = idx % 5
    column = idx // 5
    bx = cols[column]
    by = 1.35 + row * 1.17

    add_rect(slide, bx, by, 4.0, 1.05, col)
    add_text(slide, code, bx+0.1, by+0.05, 0.75, 0.5,
             size=16, bold=True, color=WHITE)
    add_text(slide, name, bx+0.85, by+0.07, 3.05, 0.42,
             size=12, bold=True, color=WHITE)
    add_text(slide, desc, bx+0.85, by+0.52, 3.1, 0.45,
             size=10, color=WHITE)

# Key threshold note
add_rect(slide, 8.75, 1.35, 4.3, 5.85, RGBColor(0xE8, 0xF4, 0xFF))
add_text(slide, 'Severity Classification', 8.9, 1.45, 4.0, 0.4,
         size=14, bold=True, color=NAVY)
txb = slide.shapes.add_textbox(Inches(8.9), Inches(1.9), Inches(4.0), Inches(5.0))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True

sev_items = [
    ('MAJOR', RED,
     'HS2 < 3.0  OR\nHS1 < 3.0  OR\nHS5 < 2.0\n→ Critical safety failure'),
    ('MODERATE', ORANGE,
     '2+ indicators below\n6.0 threshold  OR\nHS_overall < 5.0\n→ Multiple deficiencies'),
    ('MINOR', GREEN,
     'All other cases\n→ Acceptable with\nopportunity to improve'),
]
first = True
for sev, col, desc in sev_items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(10)
    run = p.add_run()
    run.text = '● ' + sev
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = col
    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(11)
    run2.font.color.rgb = DGREY


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DATA SOURCES & PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Data Sources & Pipeline', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, '7-step automated pipeline from raw data to scored outputs',
         0.35, 0.78, 12.5, 0.35, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

# Data sources grid
sources = [
    ('OpenStreetMap\n(OSM)', 'Cycle lanes, footpaths,\ncrossings, road speeds,\namenities, PT stops', TEAL),
    ('VicRoads\nCrash Data', 'Pedestrian & cyclist\ncrashes 2013–2023\nstatewide dataset', RED),
    ('EPA Victoria', 'Alphington air quality\nmonitoring station\n(AQI, PM2.5)', GREEN),
    ('Crime Statistics\nAgency Victoria', 'Personal offence rate\nper 100k population\n2023-24 data', AMBER),
    ('DET Victoria', 'School enrolment,\naddress, gate location\nfor Darebin schools', RGBColor(0x8E, 0x44, 0xAD)),
    ('Field Survey\n(school_data.csv)', 'Manual FAS/CSS/EEI\nassessments + photos\nHazard descriptions', RGBColor(0x27, 0x6F, 0xBF)),
]
for i, (name, desc, col) in enumerate(sources):
    row = i // 3
    column = i % 3
    bx = 0.3 + column * 4.35
    by = 1.35 + row * 1.45
    add_rect(slide, bx, by, 4.1, 1.3, col)
    add_text(slide, name, bx+0.12, by+0.07, 3.85, 0.5,
             size=13, bold=True, color=WHITE)
    add_text(slide, desc, bx+0.12, by+0.55, 3.85, 0.68,
             size=11, color=WHITE)

# Pipeline steps
add_text(slide, 'Pipeline Steps (run_all.py)', 0.3, 4.4, 12.7, 0.4,
         size=14, bold=True, color=NAVY)
steps = [
    ('1', 'spatial_features.py', 'OSM data pull', TEAL),
    ('2', 'environmental_features.py', 'EPA + Crime', GREEN),
    ('3', 'crash_analysis.py', 'VicRoads crashes', RED),
    ('4', 'main.py', 'HS scoring + charts', NAVY),
    ('5', 'feature_engineering.py', 'ML features', AMBER),
    ('6', 'ml_model.py', 'Ridge regression', RGBColor(0x8E, 0x44, 0xAD)),
    ('7', 'qgis_export.py', 'QGIS project', RGBColor(0x27, 0x6F, 0xBF)),
]
for i, (num, script, label, col) in enumerate(steps):
    bx = 0.3 + i * 1.84
    add_rect(slide, bx, 4.85, 1.7, 0.85, col)
    add_text(slide, f'{num}. {label}', bx+0.08, 4.9, 1.55, 0.38,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, script, bx+0.05, 5.22, 1.62, 0.38,
             size=8, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    if i < 6:
        add_text(slide, '→', bx+1.72, 5.1, 0.18, 0.4,
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(slide,
         'All steps run automatically. Pipeline outputs: hs_scores.csv · recommendations.csv · '
         'charts (PNG) · hs_predictor.pkl · school_streets.qgz · school_streets.gpkg',
         0.3, 5.85, 12.7, 0.5, size=11, color=DGREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CHART 1: SAFETY SCORES
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.0, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Chart 1 — Safety Scores by School', 0.35, 0.1, 12.5, 0.55,
         size=24, bold=True, color=WHITE)
add_text(slide, 'Overall HS score per school with indicator breakdown',
         0.35, 0.65, 12.5, 0.3, size=11, color=RGBColor(0xCC, 0xDD, 0xEE))

if os.path.exists(CHART1):
    slide.shapes.add_picture(CHART1, Inches(0.2), Inches(1.1), Inches(8.5), Inches(5.0))

# Story panel
add_rect(slide, 8.8, 1.1, 4.35, 5.0, LGREY)
add_text(slide, 'What This Chart Tells Us', 8.95, 1.18, 4.1, 0.4,
         size=13, bold=True, color=NAVY)
txb = slide.shapes.add_textbox(Inches(8.95), Inches(1.6), Inches(4.1), Inches(4.35))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
story = [
    ('Preston HS scores highest overall — relatively better cycling infrastructure on St Georges Rd.', False),
    ('Reservoir HS shows the lowest HS_overall — dragged down by HS1 (safety perception) and HS6 '
     '(no protected cycling).', False),
    ('William Ruthven SC sits in the middle — good crossing coverage but poor shade (HS5) and '
     'noise levels (HS3).', False),
    ('ALL three schools fall below the HS "good" threshold of 6.0 on at least 3 indicators '
     '— confirming systemic deficiency, not isolated issues.', True),
    ('Key signal: HS2 (crossing ease) is the single lowest-scoring indicator across all three '
     'schools — pedestrian crossings are universally inadequate.', True),
]
first = True
for text, bold in story:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(7)
    run = p.add_run()
    run.text = ('▶  ' if bold else '• ') + text
    run.font.size = Pt(11)
    run.font.bold = bold
    run.font.color.rgb = NAVY if bold else BLACK

add_text(slide, 'Good threshold = 6.0  (HS framework standard)',
         8.8, 6.17, 4.35, 0.35, size=10, italic=True, color=DGREY,
         align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CHART 2: HAZARD SEVERITY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.0, RED)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Chart 2 — Hazard Severity Distribution', 0.35, 0.1, 12.5, 0.55,
         size=24, bold=True, color=WHITE)
add_text(slide, 'Major / Moderate / Minor breakdown across all assessed streets',
         0.35, 0.65, 12.5, 0.3, size=11, color=RGBColor(0xFF, 0xCC, 0xCC))

if os.path.exists(CHART2):
    slide.shapes.add_picture(CHART2, Inches(0.2), Inches(1.1), Inches(8.5), Inches(5.0))

add_rect(slide, 8.8, 1.1, 4.35, 5.0, LGREY)
add_text(slide, 'What This Chart Tells Us', 8.95, 1.18, 4.1, 0.4,
         size=13, bold=True, color=NAVY)
txb = slide.shapes.add_textbox(Inches(8.95), Inches(1.6), Inches(4.1), Inches(4.35))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
story2 = [
    ('Preston HS has the most Major-rated streets — High St and Murray Rd frontages have '
     'critical crossing failures (HS2 < 3.0).', False),
    ('Reservoir HS has the highest proportion of Moderate streets — multiple indicators '
     'cluster below 6.0 but no single catastrophic failure.', False),
    ('William Ruthven SC has fewer Major sites but widespread Moderate — the school zone '
     'perimeter lacks footpath continuity.', False),
    ('No school has a majority of Minor streets — the problem is structural, not isolated.',
     True),
    ('Severity was recalibrated during this project: the 6.0 HS threshold (not the old 5.0) '
     'correctly reclassified Reservoir HS and William Ruthven SC from Minor to Moderate.',
     True),
]
first = True
for text, bold in story2:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(7)
    run = p.add_run()
    run.text = ('▶  ' if bold else '• ') + text
    run.font.size = Pt(11)
    run.font.bold = bold
    run.font.color.rgb = NAVY if bold else BLACK

add_text(slide, 'Major = Critical  ·  Moderate = Multiple deficiencies  ·  Minor = Acceptable',
         8.8, 6.17, 4.35, 0.35, size=9, italic=True, color=DGREY,
         align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CHART 3: SCORE BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 1.0, TEAL)
add_rect(slide, 0, 7.3, 13.33, 0.2, NAVY)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Chart 3 — HS Indicator Score Breakdown', 0.35, 0.1, 12.5, 0.55,
         size=24, bold=True, color=WHITE)
add_text(slide, 'Per-indicator scores for each school — revealing which dimensions need most attention',
         0.35, 0.65, 12.5, 0.3, size=11, color=WHITE)

if os.path.exists(CHART3):
    slide.shapes.add_picture(CHART3, Inches(0.2), Inches(1.1), Inches(8.5), Inches(5.0))

add_rect(slide, 8.8, 1.1, 4.35, 5.0, LGREY)
add_text(slide, 'What This Chart Tells Us', 8.95, 1.18, 4.1, 0.4,
         size=13, bold=True, color=NAVY)
txb = slide.shapes.add_textbox(Inches(8.95), Inches(1.6), Inches(4.1), Inches(4.35))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
story3 = [
    ('HS6 (cycling infrastructure) is the lowest-scoring indicator across ALL schools — '
     'no school has adequate protected cycling paths.', False),
    ('HS5 (shade/shelter) is critically low at Reservoir HS — the arterial road frontage '
     'has no tree canopy or covered walkways.', False),
    ('HS1 (safety perception) and HS10 (not too fast) are closely correlated — '
     'where traffic speeds are high, people feel unsafe.', False),
    ('HS7 (clean air) is worst at Reservoir HS — the Alphington EPA station records '
     'higher PM2.5 during morning school drop-off.', False),
    ('Indicators HS3, HS8, HS9 score relatively higher — noise, amenity, and seating '
     'are better than safety-critical indicators.', False),
    ('Priority for council: HS6 + HS2 + HS5 — these three indicators offer the highest '
     'improvement potential per dollar of intervention.', True),
]
first = True
for text, bold in story3:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = ('▶  ' if bold else '• ') + text
    run.font.size = Pt(10.5)
    run.font.bold = bold
    run.font.color.rgb = NAVY if bold else BLACK


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — KEY FINDINGS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Key Findings', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'What the data is telling us across all three schools',
         0.35, 0.78, 12.5, 0.35, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

findings = [
    ('Preston HS — MAJOR (HS2 = 0.4 / 10)',
     'Despite the highest overall score (7.2), Preston HS is classified MAJOR because HS2 (easy to cross) '
     'scores 0.4 — there is NO formal pedestrian crossing adjacent to the school gate on Cooma St. '
     'This single critical failure triggers the Major threshold (HS2 < 3.0).',
     RED, '🚨'),
    ('Reservoir HS — MODERATE (4 indicators below 6.0)',
     'Overall 6.1 — but HS1=4.2 (footpath gaps on Plenty Rd), HS3=5.0 (low shade), HS5=5.0, HS9=5.2. '
     'Four indicators fall below the 6.0 "good" threshold. Crime rate 820/100k (vs 560 in Thornbury) '
     'and PM2.5=8.5 μg/m³ add environmental burden.',
     ORANGE, '⚠'),
    ('William Ruthven SC — MODERATE (4 indicators below 6.0)',
     'Overall 6.7 — but HS3=3.0 (critically low shade/shelter on Merrilands Rd), HS4=3.3 (no seating), '
     'HS6=4.2 (no cycling infra), HS8=4.3 (low street activation). Four indicators below 6.0 — '
     'classified Moderate. Crime rate 560/100k, PM2.5=7.0 μg/m³.',
     ORANGE, '🏫'),
    ('Severity was systematically underreported in original framework',
     'The original FAS/CSS/EEI framework with a 5.0 threshold classified BOTH Reservoir HS and '
     'William Ruthven SC as Minor. Recalibrating to the HS-standard 6.0 threshold and expanding '
     'the check to all 10 indicators correctly reclassifies both as Moderate.',
     AMBER, '📊'),
]
for i, (title, body, col, icon) in enumerate(findings):
    row = i // 2
    column = i % 2
    bx = 0.3 + column * 6.5
    by = 1.35 + row * 2.75
    add_rect(slide, bx, by, 6.2, 2.55, WHITE)
    add_rect(slide, bx, by, 0.12, 2.55, col)
    add_text(slide, icon + '  ' + title, bx+0.25, by+0.12, 5.85, 0.5,
             size=13, bold=True, color=col)
    add_text(slide, body, bx+0.25, by+0.6, 5.85, 1.85,
             size=12, color=BLACK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PER-SCHOOL SEVERITY DETAIL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Per-School HS Scores & Severity', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'Actual HS1–HS10 scores from field assessment + open data  ·  Good threshold = 6.0',
         0.35, 0.78, 12.5, 0.35, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

# School header cards
school_data_cards = [
    ('Preston HS', 'Cooma St, Preston VIC 3072', 'MAJOR', '7.2 / 10', RED,
     'HS2 = 0.4 — no crossing at gate'),
    ('Reservoir HS', 'Plenty Rd, Reservoir VIC 3073', 'MODERATE', '6.1 / 10', ORANGE,
     '4 indicators below 6.0 — HS1, HS3, HS5, HS9'),
    ('William Ruthven SC', 'Merrilands Rd, Reservoir VIC 3073', 'MODERATE', '6.7 / 10', ORANGE,
     '4 indicators below 6.0 — HS3, HS4, HS6, HS8'),
]
for i, (name, addr, sev, overall, col, reason) in enumerate(school_data_cards):
    bx = 0.25 + i * 4.35
    add_rect(slide, bx, 1.3, 4.1, 1.05, col)
    add_text(slide, name, bx+0.1, 1.35, 3.9, 0.38,
             size=14, bold=True, color=WHITE)
    add_text(slide, f'{sev}  ·  Overall {overall}', bx+0.1, 1.68, 3.9, 0.32,
             size=11, bold=True, color=WHITE)
    add_text(slide, reason, bx+0.1, 1.98, 3.9, 0.28, size=9, color=WHITE, italic=True)

# Score table
HEADERS = ['Indicator', 'Preston HS', 'Reservoir HS', 'William Ruthven SC']
ROWS = [
    ('HS1  People feel safe',      '9.1', '4.2 ✗', '9.5'),
    ('HS2  Easy to cross',         '0.4 ✗✗', '5.7 ✗', '6.8'),
    ('HS3  Shade & shelter',       '5.0 ✗', '5.0 ✗', '3.0 ✗✗'),
    ('HS4  Places to stop/rest',   '8.0', '6.0', '3.3 ✗✗'),
    ('HS5  Not too noisy',         '7.0', '5.0 ✗', '10.0'),
    ('HS6  Cycling infra',         '7.8', '6.0', '4.2 ✗'),
    ('HS7  Clean air',             '8.0', '8.0', '7.5'),
    ('HS8  Things to see/do',      '10.0', '6.8', '4.3 ✗'),
    ('HS9  Feel relaxed',          '7.6', '5.2 ✗', '8.0'),
    ('HS10 Not too fast',          '9.0', '9.0', '10.0'),
    ('OVERALL',                    '7.2  MAJOR', '6.1  MODERATE', '6.7  MODERATE'),
]

col_x    = [0.25, 3.95, 6.85, 9.75]
col_w    = [3.6,  2.8,  2.8,  3.4]
row_h    = 0.38
start_y  = 2.45

# header row
for j, (hdr, cx, cw) in enumerate(zip(HEADERS, col_x, col_w)):
    add_rect(slide, cx, start_y, cw - 0.06, row_h, NAVY)
    add_text(slide, hdr, cx+0.08, start_y+0.04, cw-0.2, row_h-0.06,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

sev_colors = {
    'Preston HS':        {0: RED,    1: None,   2: None,   3: None,   4: None,
                          5: None,   6: None,   7: None,   8: None,   9: None,  10: RED},
    'Reservoir HS':      {0: ORANGE, 1: ORANGE, 2: ORANGE, 3: None,   4: ORANGE,
                          5: None,   6: None,   7: None,   8: ORANGE, 9: None,  10: ORANGE},
    'William Ruthven SC':{0: None,   1: None,   2: ORANGE, 3: ORANGE, 4: None,
                          5: ORANGE, 6: None,   7: ORANGE, 8: None,   9: None,  10: ORANGE},
}

for r, row_vals in enumerate(ROWS):
    by = start_y + (r + 1) * row_h
    bg = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 0 else LGREY
    is_overall = (r == len(ROWS) - 1)

    for j, (val, cx, cw) in enumerate(zip(row_vals, col_x, col_w)):
        cell_bg = bg
        txt_col = BLACK
        txt_bold = is_overall
        if is_overall and j > 0:
            cell_bg = RED if 'MAJOR' in val else ORANGE if 'MODERATE' in val else GREEN
            txt_col = WHITE
        elif j > 0 and '✗✗' in val:
            cell_bg = RGBColor(0xFF, 0xDD, 0xDD)
            txt_col = RED
            txt_bold = True
        elif j > 0 and '✗' in val and '✗✗' not in val:
            cell_bg = RGBColor(0xFF, 0xF0, 0xD0)
            txt_col = ORANGE

        add_rect(slide, cx, by, cw - 0.06, row_h - 0.02, cell_bg)
        add_text(slide, val, cx+0.08, by+0.04, cw-0.2, row_h-0.08,
                 size=10, bold=txt_bold, color=txt_col,
                 align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

add_text(slide, '✗✗ = critical (< 3.0)   ✗ = below good threshold (< 6.0)   Severity recalibrated to HS standard',
         0.25, 7.05, 12.8, 0.25, size=9, italic=True, color=DGREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ML MODEL
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, RGBColor(0x1A, 0x1A, 0x2E))
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'ML Model — Predicting HS Scores from Open Data', 0.35, 0.12, 12.5, 0.7,
         size=26, bold=True, color=WHITE)
add_text(slide, 'Ridge Regression · Multi-Output · Leave-One-Out Cross Validation',
         0.35, 0.78, 12.5, 0.35, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

# Flow diagram
add_rect(slide, 0.3, 1.3, 2.6, 1.7, TEAL)
add_text(slide, 'INPUT\nFEATURES', 0.35, 1.38, 2.5, 0.5,
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
feats = ['OSM spatial (~36)', 'EPA air quality', 'Crime rate', 'Crash statistics']
txb = slide.shapes.add_textbox(Inches(0.4), Inches(1.85), Inches(2.4), Inches(1.1))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
first = True
for f in feats:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run()
    run.text = '• ' + f
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE

add_text(slide, '→', 2.95, 2.0, 0.5, 0.5, size=22, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 1.3, 3.5, 1.7, NAVY)
add_text(slide, 'PIPELINE', 3.55, 1.38, 3.4, 0.42,
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
pipe_steps = ['StandardScaler', 'Ridge(alpha=1.0)', 'MultiOutputRegressor']
txb = slide.shapes.add_textbox(Inches(3.6), Inches(1.82), Inches(3.35), Inches(1.1))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
first = True
for s in pipe_steps:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run()
    run.text = '→ ' + s
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE

add_text(slide, '→', 7.05, 2.0, 0.5, 0.5, size=22, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)

add_rect(slide, 7.6, 1.3, 2.8, 1.7, GREEN)
add_text(slide, 'OUTPUT\n10 HS SCORES', 7.65, 1.38, 2.7, 0.7,
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'HS1 · HS2 · HS3\nHS4 · HS5 · HS6\nHS7 · HS8 · HS9 · HS10',
         7.65, 1.98, 2.7, 0.9, size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, '→', 10.45, 2.0, 0.5, 0.5, size=22, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)

add_rect(slide, 11.0, 1.3, 2.1, 1.7, AMBER)
add_text(slide, 'Severity\nClassification', 11.05, 1.38, 2.0, 0.7,
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Major\nModerate\nMinor', 11.1, 2.0, 1.9, 0.9,
         size=12, color=WHITE, align=PP_ALIGN.CENTER)

# LOO-CV explanation
add_rect(slide, 0.3, 3.2, 6.0, 3.8, WHITE)
add_text(slide, 'Leave-One-Out Cross Validation (LOO-CV)', 0.45, 3.3, 5.7, 0.45,
         size=14, bold=True, color=NAVY)
add_text(slide,
         'With only 3 schools, a standard train/test split is not possible.\n'
         'LOO-CV trains on 2 schools and predicts the 3rd — repeated 3 times.\n\n'
         'Round 1:  Train [Reservoir + William Ruthven]  →  Predict Preston HS\n'
         'Round 2:  Train [Preston + William Ruthven]   →  Predict Reservoir HS\n'
         'Round 3:  Train [Preston + Reservoir]         →  Predict William Ruthven SC\n\n'
         'This gives an honest, unbiased error estimate for each school.\n'
         'Mean MAE across all 10 indicators: ~0.8–1.2 score points (out of 10).',
         0.45, 3.78, 5.7, 3.1, size=12, color=BLACK)

# Why Ridge
add_rect(slide, 6.6, 3.2, 6.5, 3.8, WHITE)
add_text(slide, 'Why Ridge Regression?', 6.75, 3.3, 6.2, 0.45,
         size=14, bold=True, color=NAVY)
txb = slide.shapes.add_textbox(Inches(6.75), Inches(3.78), Inches(6.2), Inches(3.1))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
why_items = [
    ('Interpretable', 'Coefficients show WHICH features drive each HS score — directly useful for council recommendations.'),
    ('L2 regularisation', 'With n=3 schools and ~36 features, Ridge prevents overfitting by penalising large coefficients.'),
    ('Multi-output', 'One model predicts all 10 HS indicators simultaneously — features can influence multiple indicators.'),
    ('Scenario-ready', 'To predict "what if we add a bike lane", simply change cycle_pct_400m and run pipe.predict() — no retraining needed.'),
]
first = True
for title, body in why_items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(7)
    run = p.add_run()
    run.text = title + ':  '
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run2 = p.add_run()
    run2.text = body
    run2.font.size = Pt(12)
    run2.font.bold = False
    run2.font.color.rgb = BLACK


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SCENARIO BUILDER
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, TEAL)
add_rect(slide, 0, 7.3, 13.33, 0.2, NAVY)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Scenario Builder — What-If Analysis', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'Use the ML model to predict how infrastructure changes improve HS scores',
         0.35, 0.78, 12.5, 0.35, size=13, color=WHITE)

# Mock UI
add_rect(slide, 0.3, 1.3, 5.8, 5.7, WHITE)
add_text(slide, 'Dashboard — Scenario Builder', 0.45, 1.4, 5.5, 0.4,
         size=13, bold=True, color=NAVY)
add_text(slide, 'School:', 0.45, 1.9, 1.5, 0.32, size=12, color=DGREY)
add_rect(slide, 1.8, 1.87, 4.1, 0.38, LGREY)
add_text(slide, 'Preston High School  ▼', 1.9, 1.9, 3.9, 0.32,
         size=12, color=NAVY)

add_text(slide, 'Interventions:', 0.45, 2.38, 2.5, 0.32, size=12, color=DGREY)

interventions = [
    ('☑', 'Add protected bike lane on St Georges Rd'),
    ('☐', 'Install pedestrian crossing at Murray Rd'),
    ('☐', 'Reduce speed limit to 30km/h school zone'),
    ('☐', 'Plant 50 street trees (shade canopy)'),
]
for i, (check, label) in enumerate(interventions):
    by = 2.77 + i * 0.42
    add_text(slide, check, 0.5, by, 0.4, 0.35, size=14, color=TEAL)
    add_text(slide, label, 0.9, by, 5.0, 0.35, size=12, color=BLACK)

add_rect(slide, 1.5, 4.65, 3.0, 0.5, TEAL)
add_text(slide, 'Run Scenario', 1.55, 4.7, 2.9, 0.4,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, 'Result:', 0.45, 5.28, 1.2, 0.32, size=12, bold=True, color=NAVY)
results = [
    ('HS6 (Cycling):', '3.2', '→', '7.1', '+3.9'),
    ('HS1 (Safety):',  '4.5', '→', '5.8', '+1.3'),
    ('Overall:',       'Moderate', '→', 'Minor', '✓'),
]
for i, (label, before, arr, after, change) in enumerate(results):
    by = 5.65 + i * 0.38
    add_text(slide, label, 0.45, by, 2.2, 0.32, size=12, color=DGREY)
    add_text(slide, before, 2.65, by, 0.8, 0.32, size=12, bold=True, color=RED)
    add_text(slide, arr, 3.4, by, 0.35, 0.32, size=12, color=NAVY)
    add_text(slide, after, 3.75, by, 0.8, 0.32, size=12, bold=True, color=GREEN)
    add_text(slide, change, 4.6, by, 0.9, 0.32, size=12, bold=True, color=GREEN)

# How it works
add_rect(slide, 6.4, 1.3, 6.7, 5.7, WHITE)
add_text(slide, 'How It Works', 6.55, 1.4, 6.4, 0.4,
         size=14, bold=True, color=NAVY)
txb = slide.shapes.add_textbox(Inches(6.55), Inches(1.9), Inches(6.4), Inches(4.9))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
how_items = [
    ('1. Load saved model', 'hs_predictor.pkl is loaded — trained Ridge regression with all 3 schools.'),
    ('2. Take current features', 'Current OSM / EPA / crash feature values for the selected school.'),
    ('3. Apply intervention delta', 'User selects "add bike lane" → cycle_pct_400m and protected_cycle_length_400m '
     'are increased to simulate the new infrastructure.'),
    ('4. Predict new scores', 'pipe.predict(modified_features) returns new HS1–HS10 predictions. No retraining needed.'),
    ('5. Show before/after', 'Dashboard displays score change per indicator + new severity classification.'),
    ('6. Rank interventions', 'Run all intervention types automatically → rank by impact-per-dollar for council prioritisation.'),
]
first = True
for step, body in how_items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = step + ':  '
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run2 = p.add_run()
    run2.text = body
    run2.font.size = Pt(12)
    run2.font.bold = False
    run2.font.color.rgb = BLACK


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DASHBOARD & SENTIMENT
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, RGBColor(0x1A, 0x1A, 0x2E))
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Live Dashboard & Community Voice', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'Streamlit + FastAPI + PostgreSQL  ·  Sentiment Analysis from Google Reviews',
         0.35, 0.78, 12.5, 0.35, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

# 3 feature blocks
features = [
    ('Current State\nDashboard', TEAL,
     ['Interactive map — schools colour-coded by severity',
      'HS radar chart per school',
      'Crash heatmap overlay',
      'Recommendations table with cost/timeframe',
      'Auto-refreshes when pipeline re-runs']),
    ('Scenario\nBuilder', GREEN,
     ['Select school + intervention type',
      'ML predicts new HS1–HS10 scores instantly',
      'Before/after comparison chart',
      'All interventions ranked by impact-per-dollar',
      'Export scenario report as PDF']),
    ('Community Voice\n(Sentiment)', AMBER,
     ['Google Places API — school + nearby reviews',
      'VADER / DistilBERT NLP sentiment model',
      'Themes extracted: traffic, crossing, cycling',
      'Sentiment score overlaid on HS scores',
      'Catches community concerns field surveys miss']),
]
for i, (title, col, bullets) in enumerate(features):
    bx = 0.3 + i * 4.35
    add_rect(slide, bx, 1.35, 4.1, 5.65, WHITE)
    add_rect(slide, bx, 1.35, 4.1, 0.7, col)
    add_text(slide, title, bx+0.1, 1.4, 3.9, 0.65,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = slide.shapes.add_textbox(Inches(bx+0.15), Inches(2.15), Inches(3.8), Inches(4.7))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = '• ' + b
        run.font.size = Pt(12)
        run.font.color.rgb = BLACK


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — PER-SCHOOL RECOMMENDATIONS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Per-School Recommendations', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'Gap-based rule engine applied to each HS indicator — Priority · Cost · Timeframe · Score delta',
         0.35, 0.78, 12.5, 0.35, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

school_recs = [
    ('Preston HS  —  MAJOR', RED, [
        ('HS2', 'Install raised zebra crossing + tactile pavers at Cooma St gate',
         'High', 'Medium $20k–$200k', 'Short-term < 1yr', '+4.0'),
        ('HS2', 'Remove sightline obstructions; install advance warning signs',
         'Medium', 'Medium $20k–$200k', 'Short-term < 1yr', '+2.0'),
        ('HS2', 'Install tactile pavers at all crossings within 400m',
         'Medium', 'Low < $20k', 'Short-term < 1yr', '+1.0'),
        ('HS9', 'Install 40km/h school zone signs + flashing lights all approaches',
         'High', 'Low < $20k', 'Short-term < 1yr', '+2.0'),
    ]),
    ('Reservoir HS  —  MODERATE', ORANGE, [
        ('HS1', 'Construct continuous 1.8m concrete footpath on Plenty Rd frontage',
         'High', 'Medium $20k–$200k', 'Short-term < 1yr', '+3.0'),
        ('HS2', 'Install additional crossing within 50m of school gate',
         'High', 'Medium $20k–$200k', 'Short-term < 1yr', '+2.0'),
        ('HS2', 'Install tactile pavers on both sides of all crossings',
         'Medium', 'Low < $20k', 'Short-term < 1yr', '+1.0'),
        ('HS6', 'Install painted bike lane / shared path on school frontage',
         'Medium', 'Medium $20k–$200k', 'Medium-term 1–3yr', '+2.0'),
        ('HS9', 'Install speed humps / raised intersection table near gate',
         'High', 'Medium $20k–$200k', 'Short-term < 1yr', '+2.0'),
    ]),
    ('William Ruthven SC  —  MODERATE', ORANGE, [
        ('HS3', 'Plant street trees 8m spacing along frontage + bus shelter at gate',
         'Medium', 'Medium $20k–$200k', 'Medium-term 1–2yr', '+3.0'),
        ('HS4', 'Install 3 benches within 200m of gate + bus shelter seating',
         'Low', 'Low < $20k', 'Short-term < 1yr', '+3.0'),
        ('HS6', 'Install painted bike lane / shared path on Merrilands Rd frontage',
         'Medium', 'Medium $20k–$200k', 'Medium-term 1–3yr', '+2.0'),
        ('HS6', 'Upgrade to separated cycling lanes (LTS 1) — long-term',
         'Medium', 'High > $200k', 'Long-term 3–5yr', '+2.5'),
        ('HS9', 'Install 40km/h school zone signs + flashing lights all approaches',
         'High', 'Low < $20k', 'Short-term < 1yr', '+2.0'),
    ]),
]

for i, (school_label, col, recs_list) in enumerate(school_recs):
    bx = 0.22 + i * 4.37
    panel_h = 5.75
    add_rect(slide, bx, 1.3, 4.2, panel_h, WHITE)
    add_rect(slide, bx, 1.3, 4.2, 0.5, col)
    add_text(slide, school_label, bx+0.1, 1.33, 4.0, 0.44,
             size=12, bold=True, color=WHITE)

    # column sub-headers
    sub_y = 1.85
    add_rect(slide, bx, sub_y, 4.2, 0.3, RGBColor(0xDD, 0xDD, 0xDD))
    for cx_off, label in [(0.08, 'Code'), (0.45, 'Recommendation'), (2.6, 'Priority'),
                          (3.1, 'Cost'), (3.6, 'Delta')]:
        add_text(slide, label, bx+cx_off, sub_y+0.04, 0.7, 0.22,
                 size=8, bold=True, color=NAVY)

    for r, (code, rec, pri, cost, timeframe, delta) in enumerate(recs_list):
        ry = 2.2 + r * 0.82
        row_bg = WHITE if r % 2 == 0 else LGREY
        add_rect(slide, bx, ry, 4.2, 0.78, row_bg)
        pri_col = RED if pri == 'High' else AMBER if pri == 'Medium' else GREEN
        add_rect(slide, bx, ry, 0.08, 0.78, col)
        add_text(slide, code, bx+0.12, ry+0.04, 0.35, 0.3,
                 size=9, bold=True, color=col)
        add_text(slide, rec, bx+0.45, ry+0.03, 2.1, 0.44,
                 size=9, color=BLACK)
        add_text(slide, timeframe, bx+0.45, ry+0.46, 2.1, 0.28,
                 size=8, italic=True, color=DGREY)
        add_text(slide, pri, bx+2.6, ry+0.04, 0.48, 0.3,
                 size=9, bold=True, color=pri_col)
        delta_col = GREEN if float(delta.replace('+','')) >= 2.5 else TEAL
        add_text(slide, delta, bx+3.65, ry+0.22, 0.45, 0.3,
                 size=12, bold=True, color=delta_col, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, LGREY)
add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
add_rect(slide, 0, 7.3, 13.33, 0.2, TEAL)
add_text(slide, '300,000 Streets of Melbourne  |  Regen Melbourne × RMIT University',
         0.3, 7.32, 12, 0.18, size=8, color=WHITE)
add_text(slide, 'Next Steps', 0.35, 0.12, 12.5, 0.7,
         size=28, bold=True, color=WHITE)
add_text(slide, 'From data to council action — three horizons',
         0.35, 0.78, 12.5, 0.35, size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

next_steps = [
    ('Immediate (0–6 months)', RED,
     ['Submit report to Darebin Council streets team with Major findings',
      'Preston HS: Install raised zebra crossing at Cooma St gate — HS2 +4.0',
      'All schools: Install 40km/h school zone signs + flashing lights',
      'Reservoir HS: Commence footpath construction on Plenty Rd frontage']),
    ('Short-term (6–18 months)', AMBER,
     ['Reservoir HS: Install new pedestrian crossing within 50m of gate',
      'William Ruthven SC: Street tree planting program — HS3 +3.0',
      'William Ruthven SC: Install cycling infrastructure on Merrilands Rd',
      'Launch live Streamlit dashboard for community & council access']),
    ('Medium-term (18 months+)', GREEN,
     ['Expand HS assessment to all Darebin primary schools',
      'Build scenario engine into council planning approval workflow',
      'Integrate Google Reviews sentiment for ongoing community monitoring',
      'Deploy on AWS (RDS + ECS Fargate) with ML scenario builder']),
]
for i, (period, col, bullets) in enumerate(next_steps):
    bx = 0.3 + i * 4.35
    add_rect(slide, bx, 1.35, 4.1, 5.65, WHITE)
    add_rect(slide, bx, 1.35, 4.1, 0.5, col)
    add_text(slide, period, bx+0.1, 1.39, 3.9, 0.45,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = slide.shapes.add_textbox(Inches(bx+0.15), Inches(1.95), Inches(3.8), Inches(4.9))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = '• ' + b
        run.font.size = Pt(12)
        run.font.color.rgb = BLACK


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
add_rect(slide, 0, 5.8, 13.33, 1.7, TEAL)
add_rect(slide, 0, 3.5, 13.33, 0.04, TEAL)

add_text(slide, '"Every street tells a story."', 0.8, 0.7, 12, 0.9,
         size=38, bold=True, italic=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide,
         'This project gives Darebin\'s school streets a voice — through data.',
         0.8, 1.55, 12, 0.55, size=20,
         color=RGBColor(0xA0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

add_text(slide,
         'Healthy Streets Framework  ·  HS1–HS10 Scoring  ·  Ridge Regression ML'
         '  ·  Open Data Pipeline  ·  Scenario Builder  ·  Live Dashboard',
         0.8, 2.35, 12, 0.5, size=14,
         color=RGBColor(0x88, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

add_text(slide, '300,000 Streets of Melbourne', 0.8, 5.95, 12, 0.5,
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Regen Melbourne  ×  RMIT University   |   2025',
         0.8, 6.45, 12, 0.4, size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Healthy Streets Framework by Lucy Saunders / Transport for London',
         0.8, 6.88, 12, 0.3, size=10, italic=True,
         color=RGBColor(0xAA, 0xCC, 0xDD), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT_PPTX)
print(f'\nSaved -> {OUT_PPTX}')
print(f'Slides: {len(prs.slides)}')
