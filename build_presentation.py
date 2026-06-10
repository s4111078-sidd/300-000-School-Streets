"""
Build the School Streets Safety presentation for mixed audience:
data science, social science, Vic Health.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── colour palette ────────────────────────────────────────────────────────────
DARK_GREEN  = RGBColor(0x1B, 0x5E, 0x20)   # headings
MID_GREEN   = RGBColor(0x2E, 0x7D, 0x32)   # accent
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)   # bg tint
TEAL        = RGBColor(0x00, 0x69, 0x6E)   # callout boxes
AMBER       = RGBColor(0xF5, 0x7F, 0x17)   # warning / major
RED         = RGBColor(0xC6, 0x28, 0x28)   # critical
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL    = RGBColor(0x21, 0x21, 0x21)
LIGHT_GREY  = RGBColor(0xF5, 0xF5, 0xF5)
SLATE       = RGBColor(0x37, 0x47, 0x4F)

OUTPUTS = os.path.join(os.path.dirname(__file__), "outputs")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, left, top, width, height,
             bold=False, size=18, color=CHARCOAL, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return tb


def add_image(slide, path, left, top, width, height=None):
    if not os.path.exists(path):
        return None
    if height:
        return slide.shapes.add_picture(
            path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    return slide.shapes.add_picture(
        path, Inches(left), Inches(top), Inches(width)
    )


def header_bar(slide, title, subtitle=None):
    """Dark green header bar at top of slide."""
    add_rect(slide, 0, 0, 13.33, 1.1, DARK_GREEN)
    add_text(slide, title, 0.3, 0.08, 12.7, 0.65,
             bold=True, size=26, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.7, 12.7, 0.38,
                 size=14, color=RGBColor(0xC8, 0xE6, 0xC9))


def footer(slide, page_num):
    """Light footer strip."""
    add_rect(slide, 0, 7.1, 13.33, 0.4, LIGHT_GREEN)
    add_text(slide, "300,000 Streets of Melbourne — Regen Melbourne × RMIT University",
             0.3, 7.12, 10, 0.3, size=9, color=SLATE)
    add_text(slide, str(page_num), 12.7, 7.12, 0.5, 0.3,
             size=9, color=SLATE, align=PP_ALIGN.RIGHT)


def bullet_box(slide, bullets, left, top, width, height,
               bg=LIGHT_GREY, title=None, title_color=DARK_GREEN,
               bullet_size=13, title_size=14):
    add_rect(slide, left, top, width, height, bg)
    y = top + 0.1
    if title:
        add_text(slide, title, left + 0.15, y, width - 0.3, 0.35,
                 bold=True, size=title_size, color=title_color)
        y += 0.35
    for b in bullets:
        add_text(slide, f"  • {b}", left + 0.1, y, width - 0.2, 0.32,
                 size=bullet_size, color=CHARCOAL)
        y += 0.32


def stat_card(slide, value, label, left, top, width=1.9, height=1.1,
              bg=MID_GREEN, val_color=WHITE, lbl_color=RGBColor(0xC8, 0xE6, 0xC9)):
    add_rect(slide, left, top, width, height, bg)
    add_text(slide, value, left, top + 0.08, width, 0.55,
             bold=True, size=30, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + 0.62, width, 0.45,
             size=11, color=lbl_color, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)

# Full dark-green top band
add_rect(slide, 0, 0, 13.33, 3.6, DARK_GREEN)

# Decorative accent line
add_rect(slide, 0, 3.6, 13.33, 0.08, AMBER)

# Title
add_text(slide, "300,000 Streets of Melbourne",
         0.6, 0.45, 12.1, 1.1, bold=True, size=42, color=WHITE)
add_text(slide, "School Streets Safety Analysis — City of Darebin Pilot",
         0.6, 1.55, 12.1, 0.75, bold=False, size=26,
         color=RGBColor(0xC8, 0xE6, 0xC9))

# Org line
add_text(slide, "Regen Melbourne  ×  RMIT University",
         0.6, 2.35, 12.1, 0.6, size=18,
         color=RGBColor(0xA5, 0xD6, 0xA7))

# Subtitle band
add_rect(slide, 0, 3.68, 13.33, 2.9, LIGHT_GREEN)

# Key numbers strip
stat_card(slide, "3",      "Secondary schools\nassessed",       0.6,  3.85, bg=MID_GREEN)
stat_card(slide, "7,948",  "Vic ped/cyc crashes\n2021–2025",    2.65, 3.85, bg=TEAL)
stat_card(slide, "21",     "Darebin LGA crashes\n(2024 only)",  4.70, 3.85, bg=TEAL)
stat_card(slide, "r = 0.84","Equity–safety\ncorrelation",       6.75, 3.85, bg=AMBER)
stat_card(slide, "17:00",  "Peak crash hour\n(school pickup)",  8.80, 3.85, bg=RED)

add_text(slide, "Healthy Streets framework  |  Open data + field observation  |  ML scoring  |  Equity analysis",
         0.6, 5.2, 12.1, 0.45, size=13, color=SLATE, align=PP_ALIGN.CENTER)

add_text(slide, "June 2026", 0.6, 5.75, 12.1, 0.4,
         size=12, color=SLATE, align=PP_ALIGN.CENTER, italic=True)

footer(slide, 1)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem Statement
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "Problem Statement",
           "Why do school streets in Melbourne need urgent attention?")
add_rect(slide, 0, 1.1, 13.33, 6.0, LIGHT_GREEN)

# Left column — the human problem
add_rect(slide, 0.3, 1.3, 6.1, 5.4, WHITE)
add_text(slide, "The Human Cost", 0.5, 1.35, 5.7, 0.45,
         bold=True, size=17, color=DARK_GREEN)

problems = [
    "Darebin LGA ped/cyc crashes nearly TRIPLED\n   from 25 (2021) to 73 (2024)",
    "Peak crash hour is 17:00 — school pickup time",
    "Preston HS: 15 crashes within 400 m of school gate\n   in just 4 years",
    "William Ruthven SC: 100% of crashes near the gate\n   occurred during school hours",
    "Children walk along arterial roads with no safe\n   crossings, broken footpaths, and high-speed traffic",
    "Internationally, 30–50% of school trips that could be\n   walked are driven — primarily due to perceived danger",
]
y = 1.85
for p in problems:
    add_rect(slide, 0.45, y, 5.8, 0.58, LIGHT_GREEN)
    add_text(slide, f"• {p}", 0.55, y + 0.04, 5.6, 0.52,
             size=11.5, color=CHARCOAL)
    y += 0.65

# Right column — why this matters to health
add_rect(slide, 6.75, 1.3, 6.2, 2.55, RGBColor(0xE3, 0xF2, 0xFD))
add_text(slide, "Why It Matters — Health Lens", 6.9, 1.35, 5.9, 0.45,
         bold=True, size=17, color=RGBColor(0x01, 0x57, 0x9B))

health_bullets = [
    "Active travel to school is associated with better\n   physical and mental health outcomes",
    "Children in disadvantaged areas walk more — but\n   face worse street conditions (equity gap)",
    "Sedentary travel habits formed in childhood persist\n   into adulthood (WHO Active Transport evidence)",
    "Air quality, noise, and stress near arterial roads\n   directly impact child development",
]
yh = 1.9
for b in health_bullets:
    add_text(slide, f"• {b}", 6.9, yh, 5.9, 0.5,
             size=11.5, color=CHARCOAL)
    yh += 0.52

# Right column — equity box
add_rect(slide, 6.75, 4.05, 6.2, 2.5, RGBColor(0xFF, 0xF3, 0xE0))
add_text(slide, "The Equity Problem", 6.9, 4.1, 5.9, 0.4,
         bold=True, size=17, color=AMBER)
add_text(slide,
    "Reservoir HS catchment sits in SEIFA Decile 4 "
    "(moderate-high disadvantage) — the same suburb has the "
    "lowest pedestrian safety score (HS 6.1/10) and the most "
    "indicators below standard. Safer streets cannot be a privilege "
    "of wealthier suburbs.",
    6.9, 4.55, 5.9, 1.85, size=12, color=CHARCOAL)

footer(slide, 2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — How We Approached the Problem
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "How We Approached the Problem",
           "A reproducible, open-data pipeline grounded in an international evidence base")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Framework box
add_rect(slide, 0.3, 1.25, 4.5, 5.85, WHITE)
add_text(slide, "Healthy Streets Framework", 0.45, 1.3, 4.2, 0.4,
         bold=True, size=15, color=DARK_GREEN)
add_text(slide, "(Lucy Saunders / Transport for London)",
         0.45, 1.72, 4.2, 0.3, size=10, color=SLATE, italic=True)

indicators = [
    ("HS1", "Pedestrians from all walks of life"),
    ("HS2", "Easy to cross"),
    ("HS3", "Shade and shelter"),
    ("HS4", "Places to stop and rest"),
    ("HS5", "Not too noisy"),
    ("HS6", "Active travel (walk/cycle/PT)"),
    ("HS7", "People feel safe"),
    ("HS8", "Things to see and do"),
    ("HS9", "People feel relaxed"),
    ("HS10","Clean air"),
]
yi = 2.1
for code, name in indicators:
    add_rect(slide, 0.4, yi, 0.65, 0.3, MID_GREEN)
    add_text(slide, code, 0.4, yi + 0.02, 0.65, 0.28,
             bold=True, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, name, 1.1, yi + 0.02, 3.5, 0.28, size=10.5, color=CHARCOAL)
    yi += 0.34

# Pipeline steps
add_rect(slide, 5.1, 1.25, 7.9, 5.85, WHITE)
add_text(slide, "Our 10-Step Analysis Pipeline", 5.25, 1.3, 7.6, 0.4,
         bold=True, size=15, color=DARK_GREEN)

steps = [
    ("1", "VicRoads crash data",   "7,948 ped/cyc crashes (2021–2025), nearest school assigned"),
    ("2", "OSM spatial features",  "53 features per school at 200m / 400m / 800m buffers"),
    ("3", "Environmental data",    "EPA PM2.5 air quality + Crime Statistics Agency crime rate"),
    ("4", "HS scoring + maps",     "10-indicator scores, hazard severity, interactive maps"),
    ("5", "ML feature matrix",     "26 open-data features × 10 HS score targets"),
    ("6", "Ridge regression ML",   "LOO-CV prediction of HS scores from open data (MAE 2.88)"),
    ("7", "SEIFA analysis",        "ABS 2021 disadvantage deciles per school catchment"),
    ("8", "Equity overlay",        "SEIFA × HS safety scores — r = 0.84 correlation"),
    ("9", "Crash trend analysis",  "Year-on-year trends, school-hours breakdown, time-of-day"),
    ("10","Scenario analysis",     "What-if modelling: 10 interventions × 3 schools = 30 scenarios"),
]

ys = 1.75
for num, title, desc in steps:
    add_rect(slide, 5.2, ys, 0.42, 0.3, MID_GREEN)
    add_text(slide, num, 5.2, ys + 0.02, 0.42, 0.28,
             bold=True, size=9, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, 5.68, ys + 0.01, 2.2, 0.28,
             bold=True, size=10, color=DARK_GREEN)
    add_text(slide, desc, 7.92, ys + 0.01, 4.8, 0.28,
             size=9.5, color=CHARCOAL)
    ys += 0.34

# Data sources footnote
add_text(slide,
    "Data sources: data.vic.gov.au  |  OpenStreetMap  |  EPA Victoria AirWatch  |  "
    "Crime Statistics Agency Victoria  |  ABS Census 2021  |  ABS SEIFA 2021",
    5.15, 6.95, 7.9, 0.3, size=8.5, color=SLATE, italic=True)

footer(slide, 3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — HS Scores at a Glance (chart)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "What the Data Tells Us — Safety Scores",
           "Healthy Streets scores (0–10) across 10 indicators for 3 Darebin secondary schools")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Score table
schools = [
    ("Reservoir HS",      [4.2, 5.7, 5.0, 6.0, 5.0, 6.0, 8.0, 6.8, 5.2, 9.0], 6.1, "Moderate", AMBER),
    ("William Ruthven SC",[9.5, 6.8, 3.0, 3.3,10.0, 4.2, 7.5, 4.3, 8.0,10.0], 6.7, "Moderate", AMBER),
    ("Preston HS",        [9.1, 0.4, 5.0, 8.0, 7.0, 7.8, 8.0,10.0, 7.6, 9.0], 7.2, "Major",    RED),
]
headers = ["HS1","HS2","HS3","HS4","HS5","HS6","HS7","HS8","HS9","HS10","Overall","Severity"]

# Header row
col_w = [2.5, 0.65, 0.65, 0.65, 0.65, 0.65, 0.65, 0.65, 0.65, 0.65, 0.65, 0.75, 0.8]
x_pos = [0.3]
for w in col_w:
    x_pos.append(x_pos[-1] + w)

add_rect(slide, 0.3, 1.25, 12.73, 0.38, DARK_GREEN)
for i, h in enumerate(["School"] + headers):
    add_text(slide, h, x_pos[i] + 0.04, 1.28, col_w[i] - 0.08, 0.3,
             bold=True, size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

row_y = 1.65
for sname, scores, overall, sev, sev_col in schools:
    row_bg = WHITE if schools.index((sname, scores, overall, sev, sev_col)) % 2 == 0 else LIGHT_GREEN
    add_rect(slide, 0.3, row_y, 12.73, 0.5, row_bg)
    add_text(slide, sname, x_pos[0] + 0.05, row_y + 0.1, col_w[0] - 0.1, 0.35,
             bold=True, size=11, color=DARK_GREEN)
    for j, sc in enumerate(scores):
        cell_color = RED if sc < 3 else (AMBER if sc < 6 else MID_GREEN)
        add_rect(slide, x_pos[j+1] + 0.05, row_y + 0.09, col_w[j+1] - 0.1, 0.32, cell_color)
        add_text(slide, str(sc), x_pos[j+1] + 0.05, row_y + 0.1, col_w[j+1] - 0.1, 0.3,
                 bold=True, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    # overall  (column index 11 in col_w / x_pos)
    add_text(slide, str(overall), x_pos[11] + 0.05, row_y + 0.1, col_w[11] - 0.1, 0.3,
             bold=True, size=12, color=sev_col, align=PP_ALIGN.CENTER)
    # severity badge (column index 12 in col_w / x_pos)
    add_rect(slide, x_pos[12] + 0.03, row_y + 0.09, col_w[12] - 0.06, 0.32, sev_col)
    add_text(slide, sev, x_pos[12] + 0.03, row_y + 0.1, col_w[12] - 0.06, 0.3,
             bold=True, size=9, color=WHITE, align=PP_ALIGN.CENTER)
    row_y += 0.52

# Legend
add_rect(slide, 0.3, 2.95, 12.73, 0.3, WHITE)
add_text(slide, "Score colour key:", 0.5, 2.97, 1.5, 0.25, bold=True, size=9.5, color=CHARCOAL)
for col, lbl, xp in [(RED, "< 3.0  Critical", 1.9), (AMBER, "3.0–5.9  Below standard", 3.5), (MID_GREEN, "≥ 6.0  Meeting standard", 5.9)]:
    add_rect(slide, xp, 3.0, 0.25, 0.2, col)
    add_text(slide, lbl, xp + 0.3, 2.97, 2.2, 0.25, size=9.5, color=CHARCOAL)

# Chart image
add_image(slide, os.path.join(OUTPUTS, "chart3_score_breakdown.png"),
          0.3, 3.35, 12.7, 3.45)

footer(slide, 4)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Crash Trends
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "What the Data Tells Us — Crash Trends",
           "Victorian pedestrian & cyclist crash data 2021–2025, Darebin LGA")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Stat cards row
stat_card(slide, "192%", "increase in Darebin\nped/cyc crashes\n2021→2024",  0.3, 1.25, 2.3, 1.3, bg=RED)
stat_card(slide, "15",   "crashes near\nPreston HS gate\n(400 m, 4 yrs)",   2.8, 1.25, 2.3, 1.3, bg=AMBER)
stat_card(slide, "17:00","peak crash hour\nacross Darebin LGA",              5.3, 1.25, 2.3, 1.3, bg=TEAL)
stat_card(slide, "100%", "of William Ruthven\ncrashes during\nschool hours", 7.8, 1.25, 2.3, 1.3, bg=TEAL)
stat_card(slide, "27%",  "of Preston HS\ncrashes in\nschool hours",         10.3, 1.25, 2.3, 1.3, bg=MID_GREEN)

# Crash trend chart
add_image(slide, os.path.join(OUTPUTS, "chart_crash_trends.png"),
          0.3, 2.7, 12.7, 4.15)

footer(slide, 5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Equity Finding
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "What the Data Tells Us — Equity",
           "Socioeconomic disadvantage and pedestrian safety are strongly correlated")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Big correlation callout
add_rect(slide, 0.3, 1.25, 4.1, 2.2, RGBColor(0xFF, 0xF3, 0xE0))
add_text(slide, "r = 0.84", 0.4, 1.3, 3.9, 1.0,
         bold=True, size=54, color=AMBER, align=PP_ALIGN.CENTER)
add_text(slide, "Pearson correlation between\nSEIFA disadvantage decile\nand HS overall score",
         0.4, 2.25, 3.9, 0.9, size=13, color=CHARCOAL, align=PP_ALIGN.CENTER)

# SEIFA table
seifa_data = [
    ("Reservoir HS",      "4  (most\ndisadvantaged)", 6.1, "Moderate-high", RED),
    ("William Ruthven SC","4  (most\ndisadvantaged)", 6.7, "Moderate-high", RED),
    ("Preston HS",        "6  (moderate)",            7.2, "Moderate",      AMBER),
]
add_rect(slide, 0.3, 3.65, 4.1, 0.38, DARK_GREEN)
for i, h in enumerate(["School", "SEIFA Decile", "HS Score", "Disadvantage"]):
    add_text(slide, h, [0.35, 1.55, 2.75, 3.35][i], 3.68, [1.15, 1.15, 0.6, 0.75][i], 0.28,
             bold=True, size=9, color=WHITE)
ry = 4.05
for sn, dec, hss, dlvl, dc in seifa_data:
    add_rect(slide, 0.3, ry, 4.1, 0.5, WHITE if seifa_data.index((sn,dec,hss,dlvl,dc))%2==0 else LIGHT_GREEN)
    add_text(slide, sn, 0.35, ry+0.08, 1.15, 0.4, bold=True, size=9, color=DARK_GREEN)
    add_text(slide, dec, 1.55, ry+0.08, 1.15, 0.4, size=10, color=dc, bold=True)
    add_text(slide, str(hss), 2.75, ry+0.1, 0.55, 0.32,
             bold=True, size=14, color=dc, align=PP_ALIGN.CENTER)
    add_text(slide, dlvl, 3.35, ry+0.1, 1.0, 0.32, size=9, color=dc)
    ry += 0.52

# Equity chart
add_image(slide, os.path.join(OUTPUTS, "chart_equity_seifa.png"),
          4.6, 1.25, 8.45, 5.55)

# Key implication box
add_rect(slide, 0.3, 5.35, 4.1, 1.35, RGBColor(0xE8, 0xF5, 0xE9))
add_text(slide, "Key Implication", 0.45, 5.38, 3.8, 0.35,
         bold=True, size=13, color=DARK_GREEN)
add_text(slide,
    "The schools serving the most disadvantaged communities "
    "also have the worst safety infrastructure. This is not coincidence — "
    "it reflects decades of under-investment. Prioritised funding "
    "is both a safety and an equity imperative.",
    0.45, 5.74, 3.8, 0.9, size=10.5, color=CHARCOAL)

footer(slide, 6)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — ML Findings
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "What the Data Tells Us — Machine Learning",
           "Can open data predict pedestrian safety? Which indicators still need field surveys?")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Left — explanation
add_rect(slide, 0.3, 1.25, 5.0, 5.85, WHITE)
add_text(slide, "What We Modelled", 0.5, 1.3, 4.7, 0.4,
         bold=True, size=16, color=DARK_GREEN)
add_text(slide,
    "We trained a Ridge regression model to predict each HS indicator "
    "score from 26 open-data features (OSM, AQI, crime statistics, crash counts).\n\n"
    "Evaluation: Leave-One-Out cross-validation (n=3 schools).\n\n"
    "Mean MAE: 2.88 / 10  — illustrative with n=3; the approach scales.",
    0.5, 1.75, 4.7, 1.65, size=12, color=CHARCOAL)

add_text(slide, "What This Means in Practice", 0.5, 3.45, 4.7, 0.38,
         bold=True, size=15, color=DARK_GREEN)

insights = [
    ("Automatable (no survey needed)",
     ["HS7 — Safety (MAE 0.54): crime stats are a reliable proxy",
      "HS10 — Air quality (MAE 0.92): AQI maps directly"],
     MID_GREEN),
    ("Partially automatable",
     ["HS9 — Relaxed (MAE 1.72): speed/arterial OK",
      "HS3 — Shade (MAE 2.17): OSM tree data reasonable",
      "HS1 — Footpath (MAE 3.27): presence measurable, condition not"],
     AMBER),
    ("Requires field survey",
     ["HS2 — Easy to cross (MAE 4.51): quality ≠ presence",
      "HS8 — Things to do (MAE 4.72): activity quality not in OSM",
      "HS4 — Rest places (MAE 4.19): OSM bench data unreliable"],
     RED),
]
yi = 3.9
for title, items, col in insights:
    add_rect(slide, 0.35, yi, 4.85, 0.28, col)
    add_text(slide, title, 0.4, yi + 0.02, 4.8, 0.24,
             bold=True, size=10, color=WHITE)
    yi += 0.3
    for it in items:
        add_text(slide, f"  • {it}", 0.4, yi, 4.8, 0.28, size=10.5, color=CHARCOAL)
        yi += 0.28
    yi += 0.08

# Right — feature importance chart
add_image(slide, os.path.join(OUTPUTS, "chart_feature_importance.png"),
          5.6, 1.25, 7.45, 5.85)

footer(slide, 7)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Demographics
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "Community Context — Who Uses These Streets?",
           "ABS Census 2021 demographics for Reservoir and Preston catchments")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Left — key stats
add_rect(slide, 0.3, 1.25, 5.8, 5.85, WHITE)
add_text(slide, "Catchment Demographics (ABS 2021)", 0.5, 1.3, 5.5, 0.42,
         bold=True, size=16, color=DARK_GREEN)

demo = [
    ("Reservoir", "$1,541/wk", "11.1%", "7.7%", "~4,854"),
    ("Preston",   "$1,844/wk", "12.8%", "7.9%", "~3,278"),
]
headers_d = ["Suburb", "Median Income", "No-car HH", "PT to work", "School-age\nchildren"]
col_xd = [0.4, 1.45, 2.65, 3.75, 4.8]
col_wd = [1.0, 1.2,  1.1,  1.1,  1.3]

add_rect(slide, 0.35, 1.78, 5.7, 0.38, DARK_GREEN)
for i, h in enumerate(headers_d):
    add_text(slide, h, col_xd[i], 1.8, col_wd[i], 0.36,
             bold=True, size=9, color=WHITE, align=PP_ALIGN.CENTER)

yd = 2.18
for suburb, inc, nocar, pt, children in demo:
    add_rect(slide, 0.35, yd, 5.7, 0.5,
             LIGHT_GREEN if demo.index((suburb,inc,nocar,pt,children)) % 2 == 0 else WHITE)
    vals = [suburb, inc, nocar, pt, children]
    for i, v in enumerate(vals):
        add_text(slide, v, col_xd[i], yd + 0.1, col_wd[i], 0.35,
                 bold=(i == 0), size=11, color=DARK_GREEN if i == 0 else CHARCOAL,
                 align=PP_ALIGN.CENTER)
    yd += 0.52

add_text(slide,
    "Both catchments are lower-income relative to metropolitan Melbourne median ($1,950/week). "
    "Car ownership is below average, meaning residents are more dependent on walking, "
    "cycling, and public transport — making street safety outcomes disproportionately impactful.",
    0.5, 3.35, 5.5, 1.1, size=12, color=CHARCOAL)

add_rect(slide, 0.35, 4.55, 5.7, 2.3, RGBColor(0xE3, 0xF2, 0xFD))
add_text(slide, "Active Travel Dependency",
         0.5, 4.58, 5.5, 0.38, bold=True, size=14, color=RGBColor(0x01, 0x57, 0x9B))
add_text(slide,
    "Lower car ownership in these suburbs means that children, "
    "older residents, and lower-income households have few alternatives "
    "to walking or cycling. Dangerous streets don't just create inconvenience "
    "— they restrict access to education, employment, and essential services. "
    "This is a health equity issue, not just a traffic engineering one.",
    0.5, 5.0, 5.5, 1.75, size=11.5, color=CHARCOAL)

# Demographics chart
add_image(slide, os.path.join(OUTPUTS, "chart4_demographics.png"),
          6.35, 1.25, 6.7, 5.85)

footer(slide, 8)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — What Can Be Done (Scenario Analysis)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "What Can Be Done — Scenario Analysis",
           "Modelled impact of 10 physical interventions on Healthy Streets scores")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Explanation
add_rect(slide, 0.3, 1.25, 5.5, 2.7, WHITE)
add_text(slide, "The Scenario Engine", 0.5, 1.3, 5.2, 0.4,
         bold=True, size=16, color=DARK_GREEN)
add_text(slide,
    "Using the trained Ridge regression model and a delta method, "
    "we modelled what happens to HS scores when a specific intervention "
    "is applied (e.g. install a pedestrian crossing, plant street trees, "
    "reduce speed limit to 40 km/h).\n\n"
    "Actual field-observed scores are held as the baseline — the model "
    "contributes only the predicted change. 30 scenarios computed: "
    "10 interventions × 3 schools.",
    0.5, 1.75, 5.2, 1.95, size=12, color=CHARCOAL)

# Interventions table
add_rect(slide, 0.3, 4.05, 5.5, 0.35, DARK_GREEN)
add_text(slide, "Intervention", 0.45, 4.08, 2.2, 0.28, bold=True, size=9.5, color=WHITE)
add_text(slide, "Target", 2.7, 4.08, 0.9, 0.28, bold=True, size=9.5, color=WHITE)
add_text(slide, "Cost", 3.65, 4.08, 1.2, 0.28, bold=True, size=9.5, color=WHITE)
add_text(slide, "Lead time", 4.9, 4.08, 0.75, 0.28, bold=True, size=9.5, color=WHITE)

interventions = [
    ("Signalised pedestrian crossing",   "HS2", "$80k–$200k",     "< 1 yr"),
    ("Speed limit → 40 km/h zone",       "HS5", "$5k–$20k",       "< 6 mo"),
    ("Traffic calming (speed humps)",     "HS9", "$50k–$150k",     "< 1 yr"),
    ("Continuous footpath",              "HS1", "$100k–$300k",    "< 1 yr"),
    ("Protected bike lane",              "HS6", "$500k–$1.5M/km", "1–3 yr"),
    ("Street trees",                     "HS3", "$20k–$60k",      "< 1 yr"),
    ("Public transport stop upgrade",    "HS6", "$50k–$200k",     "1–2 yr"),
    ("Street benches / seating",         "HS4", "$5k–$15k",       "< 6 mo"),
    ("Covered shelter",                  "HS3", "$15k–$40k",      "< 6 mo"),
    ("Remove/reroute arterial traffic",  "HS5", "High / long",    "2–5 yr"),
]
iy = 4.42
for i, (name, target, cost, lead) in enumerate(interventions):
    bg = WHITE if i % 2 == 0 else LIGHT_GREEN
    add_rect(slide, 0.3, iy, 5.5, 0.32, bg)
    add_text(slide, name,   0.45, iy+0.04, 2.2, 0.25, size=9.5, color=CHARCOAL)
    add_text(slide, target, 2.7,  iy+0.04, 0.9, 0.25, bold=True, size=10, color=MID_GREEN)
    add_text(slide, cost,   3.65, iy+0.04, 1.2, 0.25, size=9, color=CHARCOAL)
    add_text(slide, lead,   4.9,  iy+0.04, 0.75,0.25, size=9, color=SLATE)
    iy += 0.32

# Preston HS scenario chart
add_image(slide,
    os.path.join(OUTPUTS, "scenario_Preston_HS_pedestrian_crossing_speed_reduction.png"),
    6.0, 1.25, 7.1, 3.3)

# Ranking chart
add_image(slide,
    os.path.join(OUTPUTS, "scenario_ranking_Reservoir_HS.png"),
    6.0, 4.65, 7.1, 2.35)

add_text(slide, "Preston HS — pedestrian crossing + speed reduction scenario",
         6.05, 4.55, 7.0, 0.28, size=9, color=SLATE, italic=True)

footer(slide, 9)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Recommendations
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "Recommendations",
           "Prioritised interventions for Darebin secondary school streets")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Per-school recommendation columns
rec_data = {
    "Preston HS\n(Major — HS2=0.4)": {
        "color": RED,
        "items": [
            ("CRITICAL", "Install signalised crossing at school gate — no safe crossing exists (HS2=0.4)", "< 1 yr", "$80k–$200k"),
            ("High",     "Reduce speed limit to 40 km/h school zone", "< 6 mo", "$5k–$20k"),
            ("High",     "Install traffic calming on Cooma St approaches", "< 1 yr", "$50k–$150k"),
            ("Medium",   "Add PT stop improvement / active travel signage", "1–2 yr", "$50k–$200k"),
        ]
    },
    "Reservoir HS\n(Moderate — 5 indicators below 6)": {
        "color": AMBER,
        "items": [
            ("High",   "Construct continuous footpath on Plenty Rd (HS1=4.2)", "< 1 yr", "$100k–$300k"),
            ("High",   "Install tactile pavers and kerb ramps at all crossings", "< 1 yr", "$20k–$80k"),
            ("High",   "Traffic calming — multi-lane road near gate (HS9=5.2)", "< 1 yr", "$50k–$150k"),
            ("Medium", "Plant street trees for shade on walking routes (HS3=5.0)", "< 1 yr", "$20k–$60k"),
        ]
    },
    "William Ruthven SC\n(Moderate — shade/rest/active travel)": {
        "color": MID_GREEN,
        "items": [
            ("High",   "Plant trees and install shelters (HS3=3.0 — lowest of 3 schools)", "< 1 yr", "$35k–$100k"),
            ("High",   "Add street seating / rest points (HS4=3.3)", "< 6 mo", "$5k–$20k"),
            ("Medium", "Improve cycling infrastructure on Merrilands Rd (HS6=4.2)", "1–3 yr", "$500k+/km"),
            ("Medium", "Add nearby amenities or improve signage to existing ones (HS8=4.3)", "1–2 yr", "$20k–$60k"),
        ]
    },
}

col_x = [0.3, 4.55, 8.8]
for ci, (school, data) in enumerate(rec_data.items()):
    cx = col_x[ci]
    add_rect(slide, cx, 1.25, 4.0, 0.55, data["color"])
    add_text(slide, school, cx + 0.1, 1.28, 3.8, 0.5,
             bold=True, size=12, color=WHITE)
    ry2 = 1.82
    for priority, rec, timeline, cost in data["items"]:
        p_col = RED if priority == "CRITICAL" else (AMBER if priority == "High" else TEAL)
        add_rect(slide, cx, ry2, 4.0, 1.15, WHITE if data["items"].index((priority,rec,timeline,cost))%2==0 else LIGHT_GREEN)
        add_rect(slide, cx + 0.05, ry2 + 0.05, 1.1, 0.25, p_col)
        add_text(slide, priority, cx + 0.05, ry2 + 0.06, 1.1, 0.22,
                 bold=True, size=8.5, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, rec, cx + 0.08, ry2 + 0.33, 3.85, 0.52, size=10, color=CHARCOAL)
        add_text(slide, f"Timeline: {timeline}  |  Cost: {cost}",
                 cx + 0.08, ry2 + 0.85, 3.85, 0.25, size=9, color=SLATE, italic=True)
        ry2 += 1.22

# Bottom CTA
add_rect(slide, 0.3, 6.6, 12.73, 0.52, DARK_GREEN)
add_text(slide,
    "All 17 recommended interventions — with priority, cost, and HS indicator — "
    "are available in the filterable recommendations table at the project web dashboard.",
    0.5, 6.63, 12.4, 0.45, size=11.5, color=WHITE, align=PP_ALIGN.CENTER)

footer(slide, 10)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Scale-up Potential & Next Steps
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK_LAYOUT)
header_bar(slide, "Scale-up Potential & Next Steps",
           "From 3 schools in Darebin to 300,000 streets across Melbourne")
add_rect(slide, 0, 1.1, 13.33, 6.4, LIGHT_GREEN)

# Three columns: scale-up, limitations, next steps
cols = [
    ("Why This Can Scale",
     MID_GREEN,
     [
         "Pipeline is fully automated — adding a new school\n  takes minutes (add gate, re-run python run_all.py)",
         "70% of HS indicators can be estimated from open data\n  alone — HS7 (crime) and HS10 (air quality) need no survey",
         "The scenario engine runs 30 what-if analyses in seconds,\n  giving planners an instant evidence base for budget decisions",
         "GIS outputs (GPKG, QGZ) are ready for council planning\n  and Transport for Victoria workflows",
         "Dashboard is static HTML — zero infrastructure,\n  shareable via GitHub Pages",
     ]),
    ("Current Limitations",
     AMBER,
     [
         "Only 3 schools assessed — ML results are illustrative,\n  not generalisable until n ≥ 20",
         "Field observations are a snapshot in time;\n  seasonal changes (shade, footpath condition) not captured",
         "OSM data quality varies by suburb — low-income areas\n  sometimes have less complete OSM coverage",
         "Scenario model predicts HS score changes, not actual\n  crash reductions — longitudinal validation needed",
         "SEIFA correlation (r=0.84) is suggestive but not causal\n  — confounders exist (land use, road type, age of suburb)",
     ]),
    ("Recommended Next Steps",
     TEAL,
     [
         "Expand to 10–20 schools across Melbourne (multiple LGAs)\n  to validate ML model generalisability",
         "Conduct longitudinal crash analysis after any intervention\n  to measure real-world impact vs scenario prediction",
         "Partner with Vic Health / DET to use SEIFA equity lens\n  for prioritising which schools are assessed first",
         "Automate annual pipeline re-run to track improvement\n  after interventions are implemented",
         "Explore integration with VicRoads ATAP and Transport\n  for Victoria active travel investment programs",
     ]),
]

for ci, (title, col, items) in enumerate(cols):
    cx = 0.3 + ci * 4.35
    add_rect(slide, cx, 1.25, 4.1, 0.45, col)
    add_text(slide, title, cx + 0.1, 1.28, 3.9, 0.38,
             bold=True, size=14, color=WHITE)
    yc = 1.75
    for item in items:
        add_rect(slide, cx, yc, 4.1, 0.88, WHITE if items.index(item)%2==0 else LIGHT_GREEN)
        add_text(slide, f"• {item}", cx + 0.1, yc + 0.08, 3.85, 0.75,
                 size=10.5, color=CHARCOAL)
        yc += 0.9

# Footer call to action
add_rect(slide, 0.3, 6.55, 12.73, 0.58, DARK_GREEN)
add_text(slide,
    "300,000 Streets of Melbourne — every street near every school in Victoria assessed, "
    "scored, and prioritised for investment  |  Regen Melbourne × RMIT University",
    0.5, 6.58, 12.4, 0.5, size=12, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

footer(slide, 11)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out_path = os.path.join(OUTPUTS, "School_Streets_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
