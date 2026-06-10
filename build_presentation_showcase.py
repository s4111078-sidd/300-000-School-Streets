"""
build_presentation_showcase.py
300,000 Streets of Melbourne — RMIT Project Showcase
Rubric-optimised: 6 min presentation + 2 min demo + 2 min Q&A
7 slides: 6 presentation + 1 demo handoff
Hishikesh Phukan | Regen Melbourne x RMIT University | 2024
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
GD   = RGBColor(0x1B, 0x5E, 0x20)
GM   = RGBColor(0x2E, 0x7D, 0x32)
GL   = RGBColor(0x43, 0xA0, 0x47)
GP   = RGBColor(0xF1, 0xF8, 0xE9)
GLIT = RGBColor(0xE8, 0xF5, 0xE9)
TEA  = RGBColor(0x00, 0x69, 0x6E)
AMB  = RGBColor(0xE6, 0x5C, 0x00)
ALT  = RGBColor(0xFF, 0xF3, 0xE0)
RED  = RGBColor(0xC6, 0x28, 0x28)
RLT  = RGBColor(0xFF, 0xEB, 0xEE)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
INK  = RGBColor(0x21, 0x21, 0x21)
SLT  = RGBColor(0x37, 0x47, 0x4F)
MUT  = RGBColor(0x78, 0x90, 0x9C)
CGRN = RGBColor(0xC8, 0xE6, 0xC9)
AGRN = RGBColor(0xA5, 0xD6, 0xA7)
DGRN = RGBColor(0x0D, 0x47, 0x18)

OUT = os.path.join(os.path.dirname(__file__), "outputs")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
BLANK = prs.slide_layouts[6]


# ── Primitives ────────────────────────────────────────────────────────────────
def R(s, l, t, w, h, fill):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def T(s, text, l, t, w, h,
      bold=False, sz=14, col=None, al=PP_ALIGN.LEFT, it=False):
    col = col or INK
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = al
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = it
    run.font.size   = Pt(sz)
    run.font.color.rgb = col

def IMG(s, fname, l, t, w, h=None):
    path = os.path.join(OUT, fname)
    if not os.path.exists(path):
        return
    if h:
        s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))


# ── Component helpers ─────────────────────────────────────────────────────────
def HDR(s, title, sub=None):
    R(s, 0, 0, 13.33, 1.05, GD)
    R(s, 0, 1.05, 13.33, 0.07, GL)
    T(s, title, 0.40, 0.08, 12.5, 0.56, bold=True, sz=22, col=WHT)
    if sub:
        T(s, sub, 0.40, 0.67, 12.5, 0.34, sz=10.5, col=AGRN)

def FTR(s, n, total=7):
    R(s, 0, 7.12, 13.33, 0.38, GLIT)
    T(s,
      "300,000 Streets of Melbourne  |  Regen Melbourne x RMIT University  |  Hishikesh Phukan",
      0.35, 7.15, 11.2, 0.28, sz=8.5, col=SLT)
    T(s, f"{n} / {total}", 12.55, 7.15, 0.60, 0.28,
      sz=8.5, col=SLT, al=PP_ALIGN.RIGHT)

def score_bg(v):
    if v < 4.0: return RED
    if v < 6.0: return AMB
    return GM


# =============================================================================
# SLIDE 1 — HOOK / TITLE
# Story beat: "Here is the one number that defines this project."
# =============================================================================
sl = prs.slides.add_slide(BLANK)

R(sl, 0, 0,    13.33, 4.24, GD)
R(sl, 0, 4.24, 13.33, 3.26, GLIT)
R(sl, 0, 4.20, 13.33, 0.08, GL)

# ── Left: project name + research question ────────────────────────────────────
T(sl, "300,000 Streets", 0.50, 0.22, 7.60, 0.88,
  bold=True, sz=42, col=WHT)
T(sl, "of Melbourne", 0.50, 1.10, 7.60, 0.64,
  bold=True, sz=36, col=AGRN)
T(sl, "School Streets Safety Analysis  |  City of Darebin", 0.50, 1.80, 7.60, 0.42,
  sz=14, col=AGRN, it=True)

R(sl, 0.50, 2.36, 7.50, 1.44, DGRN)
R(sl, 0.50, 2.36, 0.10, 1.44, GL)
T(sl,
  "Do children in disadvantaged suburbs walk to school\n"
  "on more dangerous streets than their peers?",
  0.70, 2.48, 7.18, 1.18,
  bold=True, sz=18, col=WHT)

# ── Right: THE ONE HOOK NUMBER ────────────────────────────────────────────────
T(sl, "192%", 8.30, 0.04, 4.70, 2.30,
  bold=True, sz=128, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "increase in pedestrian", 8.30, 2.34, 4.70, 0.36,
  sz=15, col=AGRN, al=PP_ALIGN.CENTER)
T(sl, "& cyclist crashes — Darebin", 8.30, 2.70, 4.70, 0.36,
  sz=15, col=AGRN, al=PP_ALIGN.CENTER)
T(sl, "2021 to 2024", 8.30, 3.12, 4.70, 0.40,
  bold=True, sz=16, col=GL, al=PP_ALIGN.CENTER)

# ── Bottom pale zone ──────────────────────────────────────────────────────────
R(sl, 0.50, 4.38, 5.50, 0.74, GD)
T(sl, "Regen Melbourne  x  RMIT University",
  0.66, 4.42, 5.20, 0.34, bold=True, sz=13, col=WHT)
T(sl, "COSC2667/COSC2777  |  2024  |  Hishikesh Phukan",
  0.66, 4.76, 5.20, 0.28, sz=10.5, col=AGRN)

mini_stats = [
    ("3",    "Schools assessed"),
    ("6",    "Open data sources"),
    ("10",   "HS indicators"),
    ("r=0.84","Equity correlation"),
]
mx = 0.50
for val, lbl in mini_stats:
    R(sl, mx, 5.24, 2.94, 0.92, WHT)
    R(sl, mx, 5.24, 0.08, 0.92, GM)
    T(sl, val, mx+0.14, 5.28, 2.72, 0.40, bold=True, sz=16, col=GD)
    T(sl, lbl, mx+0.14, 5.64, 2.72, 0.44, sz=9.5, col=SLT)
    mx += 3.10

R(sl, 0.50, 6.28, 12.33, 0.58, GD)
T(sl,
  "Every child deserves a safe walk to school.",
  0.70, 6.34, 12.00, 0.42,
  bold=True, sz=18, col=WHT, al=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — THE PROBLEM
# Story beat: "Let me make this concrete. Here are the numbers."
# =============================================================================
sl = prs.slides.add_slide(BLANK)
HDR(sl, "The Problem — A Safety Crisis at School Gates",
    "VicRoads crash data 2021-2025  |  7,948 Victorian pedestrian & cyclist crashes analysed")
R(sl, 0, 1.12, 13.33, 5.96, GLIT)

# ── Left: 3 stat blocks ───────────────────────────────────────────────────────
R(sl, 0.25, 1.22, 3.70, 1.72, RED)
R(sl, 0.25, 1.22, 3.70, 0.08, AMB)
T(sl, "192%", 0.25, 1.28, 3.70, 0.98,
  bold=True, sz=68, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "crash increase 2021 to 2024\n25 crashes  to  73 crashes in Darebin LGA",
  0.36, 2.24, 3.50, 0.62, sz=11, col=RLT, al=PP_ALIGN.CENTER)

R(sl, 0.25, 3.06, 3.70, 1.52, AMB)
R(sl, 0.25, 3.06, 3.70, 0.08, GL)
T(sl, "17:00", 0.25, 3.12, 3.70, 0.82,
  bold=True, sz=58, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "Peak crash hour across Darebin\nSchool pickup time",
  0.36, 3.92, 3.50, 0.58, sz=12, col=ALT, al=PP_ALIGN.CENTER)

R(sl, 0.25, 4.70, 3.70, 1.48, GD)
R(sl, 0.25, 4.70, 3.70, 0.08, GL)
T(sl, "100%", 0.25, 4.76, 3.70, 0.80,
  bold=True, sz=58, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "of William Ruthven SC's nearby\ncrashes during school hours",
  0.36, 5.54, 3.50, 0.56, sz=12, col=AGRN, al=PP_ALIGN.CENTER)

R(sl, 0.25, 6.28, 3.70, 0.68, DGRN)
R(sl, 0.25, 6.28, 0.09, 0.68, GL)
T(sl,
  "Equity: the most disadvantaged catchments\nhave the worst streets  (r = 0.84)",
  0.42, 6.34, 3.42, 0.56, sz=10.5, col=WHT, bold=True)

# ── Right: crash trend chart ──────────────────────────────────────────────────
IMG(sl, "chart_crash_trends.png", 4.12, 1.22, 9.00, 5.72)

R(sl, 9.50, 1.26, 2.62, 0.58, RED)
R(sl, 9.50, 1.26, 2.62, 0.08, AMB)
T(sl, "2024: 73 crashes  -- up 192%", 9.62, 1.32, 2.42, 0.44,
  sz=10.5, col=WHT, bold=True)

R(sl, 11.00, 4.10, 2.12, 0.68, AMB)
R(sl, 11.00, 4.10, 0.08, 0.68, RED)
T(sl, "Peak 17:00\nSchool pickup", 11.14, 4.16, 1.90, 0.54,
  sz=10.5, col=WHT, bold=True)

FTR(sl, 2)


# =============================================================================
# SLIDE 3 — OUR APPROACH
# Story beat: "We built a rigorous, reproducible pipeline to measure this."
# =============================================================================
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Our Approach — Healthy Streets Framework + Open Data Pipeline",
    "Lucy Saunders / Transport for London  |  10 indicators, 0-10 scale  |  6 government data sources")
R(sl, 0, 1.12, 13.33, 5.96, GLIT)

# ── Left: HS indicators as pill rows ─────────────────────────────────────────
R(sl, 0.25, 1.22, 6.00, 0.40, GD)
T(sl, "10 Healthy Streets Indicators  (target: all >= 6.0)",
  0.38, 1.26, 5.82, 0.30, bold=True, sz=12.5, col=WHT)

hs_list = [
    ("HS1",  "Footpath accessibility",      "Field survey"),
    ("HS2",  "Easy to cross",               "Field survey"),
    ("HS3",  "Shade & shelter",             "OpenStreetMap"),
    ("HS4",  "Places to rest",              "OpenStreetMap"),
    ("HS5",  "Not too noisy",               "Field survey"),
    ("HS6",  "Active travel (cycling/PT)",  "OSM + field"),
    ("HS7",  "Feel safe",                   "CSA crime rate"),
    ("HS8",  "Things to see & do",          "OpenStreetMap"),
    ("HS9",  "Feel relaxed",                "Field survey"),
    ("HS10", "Clean air (PM2.5)",           "EPA AirWatch"),
]
iy = 1.65
for code, name, source in hs_list:
    src_col = TEA if "EPA" in source else (AMB if "CSA" in source else GM)
    R(sl, 0.25, iy, 0.72, 0.36, GM)
    T(sl, code, 0.25, iy+0.04, 0.72, 0.28,
      bold=True, sz=9.5, col=WHT, al=PP_ALIGN.CENTER)
    T(sl, name, 1.04, iy+0.06, 3.14, 0.26, bold=True, sz=10, col=GD)
    T(sl, source, 4.22, iy+0.06, 1.98, 0.26, sz=9, col=src_col, it=True)
    iy += 0.39

R(sl, 0.25, 5.58, 6.00, 0.30, DGRN)
T(sl,
  "Major severity: HS2 < 4.0  OR  HS1 < 4.0  OR  HS5 < 3.0",
  0.38, 5.61, 5.80, 0.24, sz=9, col=AGRN)

# ── Right: data sources + ML ──────────────────────────────────────────────────
R(sl, 6.45, 1.22, 6.60, 0.40, TEA)
T(sl, "6 Open Data Sources", 6.60, 1.26, 6.40, 0.30,
  bold=True, sz=12.5, col=WHT)

sources = [
    (GD,  "VicRoads",                "7,948 ped/cyclist crashes  2021-2025"),
    (GD,  "OpenStreetMap",           "53 spatial features per school gate"),
    (TEA, "EPA Victoria AirWatch",   "PM2.5 air quality  --  HS10 score"),
    (AMB, "Crime Statistics Agency", "Crime rate per 100k  --  HS7 score"),
    (GM,  "ABS SEIFA 2021",          "Disadvantage decile by catchment"),
    (GM,  "ABS Census 2021",         "Car ownership, income, active travel"),
]
sy = 1.66
for bg, name, desc in sources:
    R(sl, 6.45, sy, 6.60, 0.54, WHT)
    R(sl, 6.45, sy, 0.08, 0.54, bg)
    T(sl, name, 6.62, sy+0.05, 3.00, 0.24, bold=True, sz=11, col=bg)
    T(sl, desc,  6.62, sy+0.28, 6.24, 0.22, sz=9.5, col=SLT)
    sy += 0.58

R(sl, 6.45, 5.18, 6.60, 0.46, GD)
R(sl, 6.45, 5.18, 0.09, 0.46, GL)
T(sl, "10-step automated pipeline  |  python run_all.py  |  one command",
  6.62, 5.22, 6.36, 0.36, bold=True, sz=12, col=WHT)

R(sl, 6.45, 5.72, 6.60, 0.78, DGRN)
R(sl, 6.45, 5.72, 0.09, 0.78, GL)
T(sl,
  "Machine Learning:  Ridge regression on 26 open-data features\n"
  "to predict HS scores  |  Leave-One-Out cross-validation  |  Mean MAE 2.88",
  6.62, 5.78, 6.36, 0.66, sz=10.5, col=WHT)

FTR(sl, 3)


# =============================================================================
# SLIDE 4 — KEY FINDING: THE PARADOX
# Story beat: "Our most important finding breaks intuition."
# =============================================================================
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Key Finding -- The Highest-Scoring School is the Most Dangerous",
    "Preston HS overall score = 7.2 / 10  --  yet the ONLY school rated Major severity")
R(sl, 0, 1.12, 13.33, 5.96, WHT)

# ── Color-coded score table ───────────────────────────────────────────────────
CX  = [0.18, 3.60, 6.60, 9.68]
CW  = [3.38, 2.96, 3.04, 3.47]
RH  = 0.42
HH  = 0.46
Y0  = 1.16

hs_data = [
    ("HS1",  "Footpath",      9.1,  4.2, 9.5),
    ("HS2",  "Crossings",     0.4,  5.7, 6.8),
    ("HS3",  "Shade/shelter", 5.0,  5.0, 3.0),
    ("HS4",  "Rest places",   8.0,  6.0, 3.3),
    ("HS5",  "Traffic noise", 7.0,  5.0, 10.0),
    ("HS6",  "Active travel", 7.8,  6.0, 4.2),
    ("HS7",  "Feel safe",     8.0,  8.0, 7.5),
    ("HS8",  "Activities",    10.0, 6.8, 4.3),
    ("HS9",  "Feel relaxed",  7.6,  5.2, 8.0),
    ("HS10", "Clean air",     9.0,  9.0, 10.0),
]
overall_vals = [7.2,     6.1,      6.7]
severity_lbl = ["MAJOR", "Moderate","Moderate"]
severity_col = [RED,     AMB,       AMB]

# Column headers
hdrs = [
    ("Indicator",            GD),
    ("Preston HS\nCooma St", RED),
    ("Reservoir HS\nPlenty Rd", AMB),
    ("William Ruthven SC\nMerrilands Rd", TEA),
]
for ci, ((hdr, bg), cx, cw) in enumerate(zip(hdrs, CX, CW)):
    R(sl, cx, Y0, cw, HH, bg)
    T(sl, hdr, cx+0.10, Y0+0.04, cw-0.18, HH-0.06,
      bold=True, sz=10, col=WHT, al=PP_ALIGN.CENTER)

# Indicator rows
ry = Y0 + HH
for ri, (code, name, pv, rv, wv) in enumerate(hs_data):
    is_hs2 = (code == "HS2")
    lbl_bg = ALT if is_hs2 else (GP if ri % 2 == 0 else WHT)

    # Label cell
    R(sl, CX[0], ry, CW[0], RH, lbl_bg)
    R(sl, CX[0], ry, 0.62,  RH, GM)
    T(sl, code, CX[0], ry+0.07, 0.62, RH-0.12,
      bold=True, sz=9, col=WHT, al=PP_ALIGN.CENTER)
    T(sl, name, CX[0]+0.68, ry+0.09, CW[0]-0.76, RH-0.16,
      bold=True, sz=10, col=GD)

    # Score cells
    for ci, (cx, cw, v) in enumerate(zip(CX[1:], CW[1:], [pv, rv, wv])):
        R(sl, cx, ry, cw, RH, score_bg(v))
        T(sl, f"{v:.1f}", cx, ry+0.07, cw, RH-0.12,
          bold=True, sz=14, col=WHT, al=PP_ALIGN.CENTER)

    if is_hs2:
        # Red border strip on top of HS2 row to highlight it
        R(sl, CX[1], ry, sum(CW[1:]) + 0.06, 0.04, AMB)

    ry += RH

# Divider
R(sl, 0.18, ry, 12.97, 0.04, GD)
ry += 0.06

# Overall row
for ci, (cx, cw) in enumerate(zip(CX, CW)):
    if ci == 0:
        R(sl, cx, ry, cw, HH, GD)
        T(sl, "OVERALL", cx+0.10, ry+0.10, cw-0.18, HH-0.16,
          bold=True, sz=11, col=WHT)
    else:
        v = overall_vals[ci-1]
        R(sl, cx, ry, cw, HH, score_bg(v))
        T(sl, f"{v:.1f}", cx, ry+0.08, cw, HH-0.14,
          bold=True, sz=15, col=WHT, al=PP_ALIGN.CENTER)
ry += HH + 0.04

# Severity row
for ci, (cx, cw) in enumerate(zip(CX, CW)):
    if ci == 0:
        R(sl, cx, ry, cw, 0.50, GD)
        T(sl, "SEVERITY", cx+0.10, ry+0.12, cw-0.18, 0.30,
          bold=True, sz=11, col=WHT)
    else:
        R(sl, cx, ry, cw, 0.50, severity_col[ci-1])
        T(sl, severity_lbl[ci-1], cx, ry+0.12, cw, 0.30,
          bold=True, sz=12, col=WHT, al=PP_ALIGN.CENTER)

# ── Paradox callout under Preston column ─────────────────────────────────────
R(sl, CX[1], ry+0.56, CW[1], 0.68, RED)
R(sl, CX[1], ry+0.56, CW[1], 0.07, AMB)
T(sl,
  "HS2 = 0.4 -- Crisis\nNo safe crossing at school gate\nOne failure overrides 9 green scores",
  CX[1]+0.10, ry+0.64, CW[1]-0.18, 0.52,
  bold=True, sz=9.5, col=WHT, al=PP_ALIGN.CENTER)

FTR(sl, 4)


# =============================================================================
# SLIDE 5 — EQUITY + OPEN DATA INSIGHT
# Story beat: "The pattern is not random. It follows disadvantage."
# =============================================================================
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Equity -- Disadvantaged Communities Have the Worst Streets",
    "ABS SEIFA 2021 disadvantage decile x Healthy Streets overall score  |  Pearson r = 0.84")
R(sl, 0, 1.12, 13.33, 5.96, GLIT)

# ── Left ─────────────────────────────────────────────────────────────────────
R(sl, 0.25, 1.22, 3.85, 1.58, AMB)
R(sl, 0.25, 1.22, 3.85, 0.09, RED)
T(sl, "r = 0.84", 0.25, 1.26, 3.85, 0.90,
  bold=True, sz=52, col=WHT, al=PP_ALIGN.CENTER)
T(sl, "Pearson correlation\nSEIFA decile x HS overall score",
  0.38, 2.14, 3.60, 0.56, sz=10.5, col=ALT, al=PP_ALIGN.CENTER)

# SEIFA table
R(sl, 0.25, 2.90, 3.85, 0.36, GD)
for lbl, lx in [("School", 0.32), ("Decile", 1.85), ("HS", 2.60), ("Level", 3.12)]:
    T(sl, lbl, lx, 2.93, 0.90, 0.26, bold=True, sz=8.5, col=WHT)

seifa_rows = [
    ("Reservoir HS",  "3.7", "6.1", "Mod-high", RED),
    ("W. Ruthven SC", "3.7", "6.7", "Mod-high", AMB),
    ("Preston HS",    "5.5", "7.2", "Moderate", GM),
]
ry = 3.28
for school, decile, hs, level, col in seifa_rows:
    bg = RLT if col == RED else (ALT if col == AMB else GLIT)
    R(sl, 0.25, ry, 3.85, 0.52, bg)
    T(sl, school,  0.32, ry+0.11, 1.48, 0.30, bold=True, sz=9.5, col=GD)
    T(sl, decile,  1.82, ry+0.11, 0.64, 0.30, bold=True, sz=11,  col=col, al=PP_ALIGN.CENTER)
    T(sl, hs,      2.55, ry+0.11, 0.52, 0.30, bold=True, sz=12,  col=col, al=PP_ALIGN.CENTER)
    T(sl, level,   3.12, ry+0.11, 0.90, 0.30, sz=9,      col=col)
    ry += 0.54

# Implication box
R(sl, 0.25, 4.90, 3.85, 1.10, DGRN)
R(sl, 0.25, 4.90, 0.10, 1.10, GL)
T(sl,
  "Lower income = lower car ownership\n"
  "= more walking dependency\n"
  "= more exposure to worse streets.\n\n"
  "This is an equity problem, not just a\nroad engineering problem.",
  0.44, 4.97, 3.54, 0.96, sz=10.5, col=WHT)

# Open data insight strip
R(sl, 0.25, 6.08, 3.85, 0.78, TEA)
R(sl, 0.25, 6.08, 0.10, 0.78, GL)
T(sl,
  "Open data:  HS7 (crime) and HS10 (air quality)\nfully automated -- no field survey required.",
  0.44, 6.14, 3.58, 0.66, sz=10.5, col=WHT)

# ── Right: equity chart ────────────────────────────────────────────────────────
IMG(sl, "chart_equity_seifa.png", 4.28, 1.22, 8.82, 5.72)

R(sl, 4.34, 1.26, 2.45, 0.58, RED)
R(sl, 4.34, 1.26, 0.08, 0.58, AMB)
T(sl, "Low decile + low HS score\n= priority investment zone",
  4.50, 1.30, 2.22, 0.48, sz=10, col=WHT, bold=True)

FTR(sl, 5)


# =============================================================================
# SLIDE 6 — RECOMMENDATIONS + WHAT'S NEXT
# Story beat: "Here is what needs to happen, and how it scales."
# =============================================================================
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Recommendations and What's Next",
    "Rule-based interventions per school  |  From 3 schools to 1,900 across Victoria")
R(sl, 0, 1.12, 13.33, 5.96, GLIT)

recs = [
    ("Preston HS", "MAJOR", RED, "HS2 = 0.4 -- no crossing", [
        ("CRITICAL", RED,
         "Signalised pedestrian crossing at school gate",
         "< 1 year", "$80k-$200k"),
        ("High", AMB,
         "40 km/h school zone on Cooma St",
         "< 6 months", "$5k-$20k"),
        ("High", AMB,
         "School zone signage on all approaches",
         "< 6 months", "$5k-$20k"),
    ]),
    ("Reservoir HS", "MODERATE", AMB, "HS1 = 4.2 -- broken footpath", [
        ("High", RED,
         "Continuous footpath on Plenty Rd school frontage",
         "< 1 year", "$100k-$300k"),
        ("High", AMB,
         "Traffic calming on multi-lane approach road",
         "< 1 year", "$50k-$150k"),
        ("Medium", GM,
         "Tactile pavers & kerb ramps at gate crossings",
         "< 1 year", "$20k-$80k"),
    ]),
    ("William Ruthven SC", "MODERATE", TEA, "HS3 = 3.0 -- no shade", [
        ("High", RED,
         "School zone signs on Merrilands Rd  ($5k fix)",
         "< 6 months", "$5k-$20k"),
        ("Medium", AMB,
         "Street trees -- lowest shade score of 3 schools",
         "< 1 year", "$20k-$60k"),
        ("Medium", GM,
         "Clear vegetation blocking crossing sightlines",
         "< 3 months", "< $20k"),
    ]),
]

for ci, (school, sev, scol, worst, items) in enumerate(recs):
    cx = 0.22 + ci * 4.36
    cw = 4.10

    R(sl, cx, 1.22, cw, 0.54, scol)
    R(sl, cx, 1.22, cw, 0.08, GL)
    T(sl, school, cx+0.14, 1.24, cw-0.22, 0.28,
      bold=True, sz=13, col=WHT)
    R(sl, cx+0.14, 1.54, len(sev)*0.082+0.20, 0.22, WHT)
    T(sl, sev, cx+0.16, 1.55, len(sev)*0.082+0.16, 0.18,
      bold=True, sz=8, col=scol)
    T(sl, worst, cx+0.14+len(sev)*0.082+0.28, 1.55, cw-1.10, 0.18,
      sz=9, col=WHT, bold=True)

    ry = 1.82
    for pri, pcol, rec_text, timeline, cost in items:
        ch = 1.20
        R(sl, cx, ry, cw, ch, WHT)
        R(sl, cx, ry, 0.08, ch, pcol)
        pw = len(pri)*0.082 + 0.20
        R(sl, cx+0.16, ry+0.08, pw, 0.22, pcol)
        T(sl, pri, cx+0.16, ry+0.09, pw, 0.20, bold=True, sz=8.5, col=WHT)
        T(sl, rec_text, cx+0.16, ry+0.34, cw-0.28, 0.50, sz=9.5, col=INK)
        R(sl, cx+0.14, ry+0.88, cw-0.24, 0.24, GLIT)
        T(sl, f"{timeline}  |  Cost: {cost}",
          cx+0.20, ry+0.89, cw-0.34, 0.20, sz=8.5, col=SLT, it=True)
        ry += ch + 0.06

# Scale banner
R(sl, 0.22, 5.72, 12.89, 0.56, GD)
R(sl, 0.22, 5.72, 12.89, 0.08, GL)
T(sl,
  "From 3 schools to 1,900:  Same pipeline.  Same command.  python run_all.py",
  0.42, 5.78, 12.48, 0.42,
  bold=True, sz=14, col=WHT, al=PP_ALIGN.CENTER)

R(sl, 0.22, 6.36, 12.89, 0.56, DGRN)
T(sl,
  "Every child deserves a safe walk to school.  --  300,000 Streets of Melbourne",
  0.42, 6.42, 12.48, 0.42,
  bold=True, sz=14, col=GL, al=PP_ALIGN.CENTER)

FTR(sl, 6)


# =============================================================================
# SLIDE 7 — LIVE DASHBOARD DEMO (transition slide)
# =============================================================================
sl = prs.slides.add_slide(BLANK)
HDR(sl, "Live Dashboard Demo",
    "All findings in one browser-based tool  |  Zero installation  |  GitHub Pages")
R(sl, 0, 1.12, 13.33, 5.96, WHT)

# URL banner
R(sl, 0.38, 1.22, 12.58, 0.60, GD)
R(sl, 0.38, 1.22, 12.58, 0.08, GL)
T(sl, "  s4111078-sidd.github.io/300-000-School-Streets",
  0.58, 1.28, 12.20, 0.44, bold=True, sz=19, col=WHT)

# 6 dashboard section cards (2 rows x 3 cols)
sections = [
    (GD,  "Hero Stats",
     "5 live numbers: schools, major hazards, crash count, equity r, peak hour"),
    (GM,  "Interactive Map",
     "4 layers: crash spots, walk/cycle networks, KDE heatmap, SEIFA catchments"),
    (TEA, "School Assessments",
     "Per-school: radar chart + 10-indicator bars + key hazard callout"),
    (AMB, "Scenario Explorer",
     "Pick school + intervention -- instant before/after; 30 pre-computed results"),
    (RED, "Data Analysis",
     "All 9 pipeline charts: equity, crash trends, ML results, demographics"),
    (SLT, "Recommendations",
     "Filterable table by school and priority  |  one-click CSV download"),
]
card_w = 4.04
card_h = 0.96
gap    = 0.07
for i, (bg, name, desc) in enumerate(sections):
    cx = 0.38 + (i % 3) * (card_w + gap)
    cy = 1.94 + (i // 3) * (card_h + gap)
    R(sl, cx, cy, card_w, card_h, bg)
    R(sl, cx, cy, 0.09, card_h, GL)
    T(sl, name, cx+0.18, cy+0.10, card_w-0.28, 0.28,
      bold=True, sz=12, col=WHT)
    T(sl, desc, cx+0.18, cy+0.42, card_w-0.28, 0.48,
      sz=9.5, col=CGRN)

# Demo flow strip
R(sl, 0.38, 4.08, 12.58, 0.46, DGRN)
R(sl, 0.38, 4.08, 0.10, 0.46, GL)
T(sl,
  "Demo:  Hero stats  ->  Map  ->  School Assessment tab  ->  Scenario Explorer  ->  Recommendations",
  0.58, 4.12, 12.24, 0.36, sz=11, col=WHT)

# Left: heatmap preview
IMG(sl, "heatmap.png", 0.38, 4.66, 5.68, 2.22)

# Right: 3 feature callouts
for ry_off, bg, txt in [
    (0.00, GD,  "30 pre-computed scenarios -- instant in browser"),
    (0.64, TEA, "No Python, no backend -- fully static GitHub Pages"),
    (1.28, GM,  "Recommendations CSV download -- one click"),
]:
    R(sl, 6.28, 4.66+ry_off, 6.68, 0.56, bg)
    R(sl, 6.28, 4.66+ry_off, 0.10, 0.56, GL)
    T(sl, txt, 6.48, 4.70+ry_off, 6.40, 0.42, bold=True, sz=12, col=WHT)

FTR(sl, 7)


# =============================================================================
# Save
# =============================================================================
OUT_PATH = os.path.join(OUT, "School_Streets_Showcase.pptx")
prs.save(OUT_PATH)

print()
print("=" * 62)
print("  Saved: School_Streets_Showcase.pptx")
print("  7 slides: 6 presentation + 1 demo handoff")
print()
print("  Slide timing guide (6-minute total):")
print("  1. Title / Hook          --  45 sec")
print("  2. The Problem           --  75 sec")
print("  3. Our Approach          --  60 sec")
print("  4. Key Finding (paradox) --  90 sec")
print("  5. Equity + Open Data    --  60 sec")
print("  6. Recommendations       --  60 sec")
print("  7. Demo handoff          --  (switch to browser)")
print("=" * 62)
